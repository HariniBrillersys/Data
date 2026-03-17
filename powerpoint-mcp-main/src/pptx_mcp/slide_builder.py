"""
Slide Builder - Fill existing template placeholders with content.

CORE DESIGN RULE: ALL text content goes into EXISTING template placeholders
identified by their idx. We NEVER call slide.shapes.add_textbox() or
create new text shapes. This preserves template fonts, colors, positions.

Images can be placed in three modes:
  - "fill" (default): inserted into a PICTURE placeholder, cropped to fill.
    Best for storytelling/decorative images.
  - "fit": placed as a freestanding picture shape, scaled proportionally
    with no cropping.  Best for charts, diagrams, and data images where
    the entire content must remain visible.  Multiple images are arranged
    in a clean grid layout.
  - "collage": multiple images arranged with overlapping, staggered
    positions and slight size variation for a dynamic visual feel.
    Single images behave like "fit".

The ``image`` field accepts a path string for a single image, or a list
of strings / dicts for multi-image placement (max 9 per slide).

Shape annotations (arrows, callouts, etc.) are handled separately by
shape_builder.py and are layered on top of filled placeholders.
"""

from __future__ import annotations

import copy
import logging
import math
from pathlib import Path
from typing import Union

from PIL import Image as PILImage
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

from .overflow import (
    _AVG_CHAR_WIDTH_EM,
    _EMU_PER_PT,
    _LINE_HEIGHT_FACTOR,
    _read_margins,
    resolve_font_size,
)
from .template_engine import LayoutInfo, PlaceholderInfo, PlaceholderRole, TemplateInfo
from .text_formatter import (
    BulletSpec,
    apply_bullet_type,
    copy_paragraph_properties,
    copy_run_properties,
)

log = logging.getLogger(__name__)


class SlideProxy:
    """
    Proxy for a single slide that maps semantic roles to placeholder indices.

    Provides methods like set_title(), set_content() that internally resolve
    the correct placeholder idx from the layout's role mapping.
    """

    def __init__(self, slide, layout_info: LayoutInfo):
        self._slide = slide
        self._layout = layout_info

        # Build role -> placeholder idx mapping
        self._role_map: dict[PlaceholderRole, int] = {}
        # Build idx -> placeholder info mapping for metadata access
        self._placeholder_info: dict[int, PlaceholderInfo] = {}
        for ph_info in layout_info.placeholders:
            if ph_info.role not in self._role_map:
                self._role_map[ph_info.role] = ph_info.idx
            self._placeholder_info[ph_info.idx] = ph_info

    def fill(self, content: dict) -> dict:
        """
        Fill slide placeholders from a content dict.

        Args:
            content: Dict mapping field names to content.
                Supported keys: title, subtitle, body, content, content_left,
                content_right, content_1, content_2, content_3, notice, image,
                image_mode, notes.

                ``image`` accepts a path string or a list of path strings /
                dicts for multi-image placement (max 9).  Each dict in the
                list may contain ``path`` (required) and ``z_order`` (int,
                0 = back, higher = front).

                ``image_mode`` can be ``"fill"`` (default), ``"fit"`` (grid),
                or ``"collage"`` (overlapping).

        Returns:
            Dict mapping field names to success/failure status.
        """
        results = {}

        image_mode = content.get("image_mode", "fill")

        def _set_collage(img):
            return self.set_picture_multi(img, mode="collage")

        def _set_fit(img):
            return self.set_picture_multi(img, mode="fit")

        if image_mode == "collage":
            image_setter = _set_collage
        elif image_mode == "fit":
            image_setter = _set_fit
        else:
            # "fill" mode: if a list is given, use first image only
            image_setter = self._set_picture_fill_compat

        field_map = {
            "title": self.set_title,
            "subtitle": self.set_subtitle,
            "body": self.set_body,
            "content": self.set_content,
            "content_left": self.set_content_left,
            "content_right": self.set_content_right,
            "content_1": lambda c: self.set_content_by_number(1, c),
            "content_2": lambda c: self.set_content_by_number(2, c),
            "content_3": lambda c: self.set_content_by_number(3, c),
            "notice": self.set_notice,
            "image": image_setter,
            "notes": self.set_notes,
        }

        for field_name, setter in field_map.items():
            if field_name in content:
                results[field_name] = setter(content[field_name])

        return results

    # --- Role-based setters ---

    def set_title(self, text: str) -> bool:
        idx = self._role_map.get(PlaceholderRole.TITLE)
        if idx is None:
            return False
        return self._set_text(idx, text, enable_shrink=True)

    def set_subtitle(self, text: str) -> bool:
        idx = self._role_map.get(PlaceholderRole.SUBTITLE)
        if idx is None:
            return False
        return self._set_text(idx, text, enable_shrink=True)

    def set_body(self, text: str) -> bool:
        idx = self._role_map.get(PlaceholderRole.BODY)
        if idx is None:
            return False
        return self._set_text(idx, text)

    def set_content(self, content: Union[str, list]) -> bool:
        idx = self._role_map.get(PlaceholderRole.CONTENT)
        if idx is None:
            return False
        return self._set_content(idx, content)

    def set_content_left(self, content: Union[str, list]) -> bool:
        idx = self._role_map.get(PlaceholderRole.CONTENT_LEFT)
        if idx is None:
            return False
        return self._set_content(idx, content)

    def set_content_right(self, content: Union[str, list]) -> bool:
        idx = self._role_map.get(PlaceholderRole.CONTENT_RIGHT)
        if idx is None:
            return False
        return self._set_content(idx, content)

    def set_content_by_number(self, number: int, content: Union[str, list]) -> bool:
        role_map = {
            1: PlaceholderRole.CONTENT_1,
            2: PlaceholderRole.CONTENT_2,
            3: PlaceholderRole.CONTENT_3,
        }
        role = role_map.get(number)
        if role is None:
            return False
        idx = self._role_map.get(role)
        if idx is None:
            return False
        return self._set_content(idx, content)

    def set_notice(self, content: Union[str, list]) -> bool:
        idx = self._role_map.get(PlaceholderRole.NOTICE)
        if idx is None:
            return False
        return self._set_content(idx, content)

    def set_picture(self, image_path: str) -> bool:
        """Insert an image into the picture placeholder (crop-to-fill).

        When ``image_mode`` is ``"fill"``, the image is inserted into the
        picture placeholder and cropped to fill it.  If the placeholder has
        a mask shape (crop geometry), the image integrates with the template
        design — this is the intended behavior for decorative images,
        especially on title/intro slides.
        """
        idx = self._role_map.get(PlaceholderRole.PICTURE)
        if idx is None:
            log.warning("No picture placeholder in layout '%s'", self._layout.name)
            return False

        path = Path(image_path)
        if not path.exists():
            log.warning("Image not found: %s", image_path)
            return False

        ph = self._find_placeholder(idx)
        if ph is None:
            return False

        try:
            ph.insert_picture(str(path))
            return True
        except Exception as e:
            log.warning("Could not insert picture: %s", e)
            return False

    def _set_picture_fill_compat(self, image_input: Union[str, list]) -> bool:
        """Handle ``image_mode: "fill"`` with backward compatibility.

        If *image_input* is a list, use the first image and warn.
        """
        if isinstance(image_input, list):
            if not image_input:
                return False
            first = image_input[0]
            path = first["path"] if isinstance(first, dict) else first
            log.warning("image_mode 'fill' does not support multiple images; using first image only")
            return self.set_picture(path)
        return self.set_picture(image_input)

    def set_picture_fit(self, image_path: str) -> bool:
        """Place a single image as a freestanding shape (no crop).

        Convenience wrapper kept for backward compatibility.
        """
        return self.set_picture_multi(image_path, mode="fit")

    def set_picture_multi(
        self,
        image_input: Union[str, list],
        mode: str = "fit",
    ) -> bool:
        """Place one or more images as freestanding shapes.

        Args:
            image_input: A path string for a single image, or a list of
                path strings / dicts (``{"path": str, "z_order": int}``)
                for multi-image placement.
            mode: ``"fit"`` for a clean grid layout, ``"collage"`` for
                overlapping staggered arrangement.

        Returns:
            True if at least one image was placed successfully.
        """
        # --- Single image path ---
        if isinstance(image_input, str):
            return self._place_single_image(image_input)

        # --- Multi-image list ---
        if not isinstance(image_input, list) or not image_input:
            return False

        images = self._normalize_image_list(image_input)
        if not images:
            return False

        area = self._get_image_area()

        try:
            if mode == "collage":
                return self._place_image_collage(images, area)
            else:
                return self._place_image_grid(images, area)
        except Exception as e:
            log.warning("Could not place multi-image (%s): %s", mode, e)
            return False

    # --- Single image placement (internal) ---

    def _place_single_image(self, image_path: str) -> bool:
        """Place a single image as a freestanding shape, scaled to fit."""
        path = Path(image_path)
        if not path.exists():
            log.warning("Image not found: %s", image_path)
            return False

        try:
            with PILImage.open(path) as img:
                img_w, img_h = img.size

            avail_left, avail_top, avail_w, avail_h = self._get_image_area()

            scale = min(avail_w / img_w, avail_h / img_h)
            final_w = int(img_w * scale)
            final_h = int(img_h * scale)

            left = avail_left + (avail_w - final_w) // 2
            top = avail_top + (avail_h - final_h) // 2

            self._slide.shapes.add_picture(
                str(path),
                Emu(left),
                Emu(top),
                Emu(final_w),
                Emu(final_h),
            )
            return True
        except Exception as e:
            log.warning("Could not place fitted picture: %s", e)
            return False

    # --- Multi-image helpers ---

    _MAX_IMAGES = 9

    def _normalize_image_list(self, image_input: list) -> list[tuple[Path, int]]:
        """Validate and normalize a list of image entries.

        Each entry is either a path string or a dict with ``path`` and
        optional ``z_order``.  Returns ``[(Path, z_order), ...]`` sorted
        by *z_order* ascending (lowest = placed first = back).
        """
        result: list[tuple[Path, int]] = []

        for i, entry in enumerate(image_input):
            if isinstance(entry, dict):
                raw_path = entry.get("path", "")
                z_order = entry.get("z_order", i)
            else:
                raw_path = str(entry)
                z_order = i

            path = Path(raw_path)
            if not path.exists():
                log.warning("Image not found (skipped): %s", raw_path)
                continue
            result.append((path, z_order))

        if len(result) > self._MAX_IMAGES:
            log.warning(
                "Too many images (%d); capping at %d",
                len(result),
                self._MAX_IMAGES,
            )
            result = result[: self._MAX_IMAGES]

        result.sort(key=lambda x: x[1])
        return result

    @staticmethod
    def _grid_dimensions(n: int, area_w: int, area_h: int) -> tuple[int, int]:
        """Calculate (cols, rows) for *n* images in an area."""
        if n <= 0:
            return (0, 0)
        aspect = area_w / area_h if area_h else 1.0
        cols = max(1, round(math.sqrt(n * aspect)))
        rows = math.ceil(n / cols)
        # Avoid excessive empty cells
        while cols > 1 and (rows - 1) * cols >= n:
            cols -= 1
            rows = math.ceil(n / cols)
        return cols, rows

    def _place_image_grid(
        self,
        images: list[tuple[Path, int]],
        area: tuple[int, int, int, int],
    ) -> bool:
        """Arrange images in a clean grid with gaps (no overlap)."""
        avail_left, avail_top, avail_w, avail_h = area
        n = len(images)
        cols, rows = self._grid_dimensions(n, avail_w, avail_h)

        gap_ratio = 0.02
        gap_w = int(avail_w * gap_ratio)
        gap_h = int(avail_h * gap_ratio)
        cell_w = (avail_w - gap_w * max(cols - 1, 0)) // max(cols, 1)
        cell_h = (avail_h - gap_h * max(rows - 1, 0)) // max(rows, 1)

        placed = 0
        for idx, (path, _z) in enumerate(images):
            row = idx // cols
            col = idx % cols

            # Center partial last row
            items_in_row = min(cols, n - row * cols)
            row_offset = (cols - items_in_row) * (cell_w + gap_w) // 2

            x = avail_left + row_offset + col * (cell_w + gap_w)
            y = avail_top + row * (cell_h + gap_h)

            if self._place_image_in_cell(path, x, y, cell_w, cell_h):
                placed += 1

        return placed > 0

    # Deterministic per-image variation tables for collage layout
    _SIZE_FACTORS = [1.02, 0.97, 1.04, 0.98, 1.01, 0.96, 1.03, 0.99, 1.02]
    _OFFSET_X_RATIOS = [-0.02, 0.03, -0.01, 0.02, -0.03, 0.01, -0.02, 0.03, -0.01]
    _OFFSET_Y_RATIOS = [0.02, -0.01, 0.03, -0.02, 0.01, -0.03, 0.02, -0.01, 0.03]

    def _place_image_collage(
        self,
        images: list[tuple[Path, int]],
        area: tuple[int, int, int, int],
    ) -> bool:
        """Arrange images with overlap, stagger, and size variation."""
        avail_left, avail_top, avail_w, avail_h = area
        n = len(images)
        cols, rows = self._grid_dimensions(n, avail_w, avail_h)

        # Overlap factor: cells are larger than strict grid → they overlap
        overlap = 0.12
        # Base cell size (without overlap)
        base_cell_w = avail_w // max(cols, 1)
        base_cell_h = avail_h // max(rows, 1)

        placed = 0
        for idx, (path, _z) in enumerate(images):
            row = idx // cols
            col = idx % cols

            # Center partial last row
            items_in_row = min(cols, n - row * cols)
            row_offset = (cols - items_in_row) * base_cell_w // 2

            # Base position (grid center for this cell)
            cx = avail_left + row_offset + col * base_cell_w + base_cell_w // 2
            cy = avail_top + row * base_cell_h + base_cell_h // 2

            # Per-image variation
            size_factor = self._SIZE_FACTORS[idx % len(self._SIZE_FACTORS)]
            dx = int(self._OFFSET_X_RATIOS[idx % len(self._OFFSET_X_RATIOS)] * base_cell_w)
            dy = int(self._OFFSET_Y_RATIOS[idx % len(self._OFFSET_Y_RATIOS)] * base_cell_h)

            # Effective cell with overlap
            eff_w = int(base_cell_w * (1 + overlap) * size_factor)
            eff_h = int(base_cell_h * (1 + overlap) * size_factor)

            # Top-left from center + offset
            x = cx - eff_w // 2 + dx
            y = cy - eff_h // 2 + dy

            if self._place_image_in_cell(path, x, y, eff_w, eff_h):
                placed += 1

        return placed > 0

    def _place_image_in_cell(
        self,
        path: Path,
        cell_left: int,
        cell_top: int,
        cell_w: int,
        cell_h: int,
    ) -> bool:
        """Scale an image proportionally to fit a cell and add it."""
        try:
            with PILImage.open(path) as img:
                img_w, img_h = img.size

            scale = min(cell_w / img_w, cell_h / img_h)
            final_w = int(img_w * scale)
            final_h = int(img_h * scale)

            # Center within the cell
            left = cell_left + (cell_w - final_w) // 2
            top = cell_top + (cell_h - final_h) // 2

            self._slide.shapes.add_picture(
                str(path),
                Emu(left),
                Emu(top),
                Emu(final_w),
                Emu(final_h),
            )
            return True
        except Exception as e:
            log.warning("Could not place image %s: %s", path, e)
            return False

    def set_notes(self, notes: Union[str, list[str]]) -> bool:
        """Set speaker notes for the slide.

        Args:
            notes: A string for a single note, or a list of strings for
                multiple talking points (each becomes a paragraph).

        Returns:
            True if notes were set successfully.
        """
        try:
            notes_slide = self._slide.notes_slide
            tf = notes_slide.notes_text_frame

            if isinstance(notes, str):
                tf.text = notes
            elif isinstance(notes, list) and notes:
                tf.text = notes[0]
                for point in notes[1:]:
                    p = tf.add_paragraph()
                    p.text = str(point)
            else:
                return False

            return True
        except Exception as e:
            log.warning("Could not set speaker notes: %s", e)
            return False

    # --- Internal helpers ---

    # Gap between text content and image when avoiding overlap (0.25 inch)
    _TEXT_IMAGE_GAP = 228600

    def _estimate_placeholder_text_bottom(self, placeholder) -> int | None:
        """Estimate the bottom edge of actual text content in a placeholder.

        Calculates where the text actually ends (not the full placeholder
        height) by counting lines and multiplying by line height.

        Returns:
            Bottom edge in EMU, or None if the placeholder has no meaningful
            text content.
        """
        if not hasattr(placeholder, "text_frame"):
            return None

        tf = placeholder.text_frame
        # Collect all text across paragraphs
        paragraphs = tf.paragraphs
        if not paragraphs:
            return None

        # Check if there's any actual text (not just empty template content)
        full_text = "\n".join(p.text for p in paragraphs)
        stripped = full_text.strip()
        if not stripped:
            return None

        # Get font size and margins for this placeholder
        font_pt = resolve_font_size(placeholder)
        ml, mr, mt, mb = _read_margins(placeholder)

        # Estimate usable width and lines needed
        usable_w_emu = placeholder.width - ml - mr
        usable_w_pt = usable_w_emu / _EMU_PER_PT
        char_width_pt = font_pt * _AVG_CHAR_WIDTH_EM
        chars_per_line = max(1, int(usable_w_pt / char_width_pt))

        total_lines = 0
        for p in paragraphs:
            text = p.text.strip()
            if not text:
                total_lines += 1
                continue
            total_lines += max(1, -(-len(text) // chars_per_line))

        # Convert lines to EMU height
        line_height_emu = int(font_pt * _LINE_HEIGHT_FACTOR * _EMU_PER_PT)
        text_height_emu = total_lines * line_height_emu + mt + mb

        return placeholder.top + text_height_emu

    def _adjust_area_for_text(
        self,
        area: tuple[int, int, int, int],
    ) -> tuple[int, int, int, int]:
        """Push image area below any filled text placeholders that overlap it.

        Scans all placeholders on the slide that contain actual text, estimates
        where that text ends, and shifts the image area down if it would
        otherwise overlap.  This avoids charts/images covering slide text while
        still using available space efficiently (the adjustment is based on
        actual text extent, not the full placeholder height).

        Args:
            area: (left, top, width, height) in EMU — the candidate image area.

        Returns:
            Adjusted (left, top, width, height) that avoids text overlap.
        """
        area_left, area_top, area_w, area_h = area
        area_right = area_left + area_w
        area_bottom = area_top + area_h

        max_text_bottom = area_top  # start with no adjustment

        for ph in self._slide.placeholders:
            # Skip placeholders that don't overlap horizontally
            ph_left = ph.left
            ph_right = ph.left + ph.width
            if ph_right <= area_left or ph_left >= area_right:
                continue

            text_bottom = self._estimate_placeholder_text_bottom(ph)
            if text_bottom is not None and text_bottom > max_text_bottom:
                max_text_bottom = text_bottom

        # Only adjust if text extends into the image area
        if max_text_bottom <= area_top:
            return area

        new_top = max_text_bottom + self._TEXT_IMAGE_GAP
        new_height = area_bottom - new_top

        # Safety: don't make image area unusably small (< 15% of original)
        if new_height < area_h * 0.15:
            log.debug(
                "Text overlap adjustment would leave too little image space (%d vs %d EMU); skipping adjustment.",
                new_height,
                area_h,
            )
            return area

        return area_left, new_top, area_w, new_height

    def _get_image_area(self) -> tuple[int, int, int, int]:
        """Determine the available area for a freestanding image (EMU).

        Strategy:
        1. If a PICTURE placeholder exists, use its bounds.
        2. If a CONTENT placeholder exists, use its bounds.
        3. Otherwise, use the area below the title (or full slide if no title).

        After determining the base area, adjusts for any filled text
        placeholders that overlap — pushing the image below the actual text
        content so charts don't cover slide text.

        Returns:
            (left, top, width, height) in EMU.
        """
        # Try picture placeholder bounds first
        pic_info = self._layout.get_by_role(PlaceholderRole.PICTURE)
        if pic_info is not None:
            area = pic_info.left, pic_info.top, pic_info.width, pic_info.height
            return self._adjust_area_for_text(area)

        # Try content placeholder bounds
        content_info = self._layout.get_by_role(PlaceholderRole.CONTENT)
        if content_info is not None:
            area = content_info.left, content_info.top, content_info.width, content_info.height
            return self._adjust_area_for_text(area)

        # Fallback: area below the title, with margins
        # Get slide dimensions from the presentation part
        try:
            prs_part = self._slide.part.package.presentation_part
            slide_w = prs_part.presentation.slide_width or 12192000
            slide_h = prs_part.presentation.slide_height or 6858000
        except Exception:
            # Hard-coded standard 16:9 slide dimensions as last resort
            slide_w = 12192000
            slide_h = 6858000

        margin = 457200  # 0.5 inch

        title_info = self._layout.get_by_role(PlaceholderRole.TITLE)
        if title_info is not None:
            top = title_info.top + title_info.height + margin // 2
        else:
            top = margin

        area = margin, top, slide_w - 2 * margin, slide_h - top - margin
        return self._adjust_area_for_text(area)

    def _find_placeholder(self, idx: int):
        """Find a placeholder on the slide by its idx."""
        for ph in self._slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    def _set_text(self, idx: int, text: str, enable_shrink: bool = False) -> bool:
        """Set simple text in a placeholder by idx.

        Args:
            idx: Placeholder index.
            text: Text content.
            enable_shrink: If True, enable normAutofit (shrink-to-fit) on the
                placeholder. Used for titles/subtitles to prevent wrapping.
        """
        ph = self._find_placeholder(idx)
        if ph is None:
            return False

        try:
            if hasattr(ph, "text_frame"):
                tf = ph.text_frame
                # Use the first paragraph to preserve formatting from template
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                    # Clear existing runs but keep paragraph formatting
                    for run in p.runs:
                        run.text = ""
                    if p.runs:
                        p.runs[0].text = text
                    else:
                        run = p.add_run()
                        run.text = text
                else:
                    tf.text = text

                # Enable shrink-to-fit for title/subtitle to prevent text wrapping
                if enable_shrink:
                    self._enable_autofit(ph)
            else:
                ph.text = text
            return True
        except Exception as e:
            log.warning("Could not set text in placeholder %d: %s", idx, e)
            return False

    def _set_content(self, idx: int, content: Union[str, list, dict], level: int = 0) -> bool:
        """Set content (text or bullet list) in a placeholder by idx.

        Supports:
        - Plain strings: "text"
        - Plain lists: ["item1", "item2"]
        - Extended dict format: {"items": [...], "bullet": "number"|"bullet"|"none"}
        - Nested lists: ["parent", ["child1", "child2"]]
        - Per-item dicts: [{"text": "...", "bold": True, "level": 1}]
        """
        ph = self._find_placeholder(idx)
        if ph is None:
            return False

        try:
            if not hasattr(ph, "text_frame"):
                return False

            tf = ph.text_frame

            # Detect and unwrap extended dict format: {"items": [...], "bullet": "number"}
            list_bullet_spec = None
            if isinstance(content, dict) and "items" in content:
                bullet_type = content.get("bullet", "auto")
                # Support single-character custom bullets: {"bullet": "✓"}
                if len(bullet_type) == 1 and bullet_type not in ("a",):  # single char = custom bullet
                    list_bullet_spec = BulletSpec(type="bullet", char=bullet_type)
                else:
                    list_bullet_spec = BulletSpec(
                        type=bullet_type,
                        start_at=content.get("start_at", 1),
                        scheme=content.get("scheme", ""),
                    )
                content = content["items"]  # Unwrap to list for processing below

            if isinstance(content, str):
                # Single text block - use first paragraph to keep template formatting
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                    for run in p.runs:
                        run.text = ""
                    if p.runs:
                        p.runs[0].text = content
                    else:
                        p.text = content
                    p.level = level
                else:
                    tf.text = content

            elif isinstance(content, list):
                # Flatten nested lists into (text, level, overrides, bullet_spec) tuples
                flat_items = self._flatten_content_list(content, level)

                first_pPr = None  # Captured from first paragraph for copying
                first_run_elem = None  # Captured from first paragraph's first run

                for i, (text, item_level, item_overrides, item_bullet_spec) in enumerate(flat_items):
                    if i == 0 and tf.paragraphs:
                        # Reuse first paragraph to preserve template formatting
                        p = tf.paragraphs[0]
                        # Capture pPr BEFORE modifying text (deep-copy for isolation)
                        source_pPr = p._element.find(qn("a:pPr"))
                        if source_pPr is not None:
                            first_pPr = copy.deepcopy(source_pPr)
                        # Capture first run element for run formatting propagation
                        first_r = p._element.find(qn("a:r"))
                        if first_r is not None:
                            first_run_elem = copy.deepcopy(first_r)

                        # Clear existing text in runs (preserve formatting)
                        for run in p.runs:
                            run.text = ""
                        if p.runs:
                            p.runs[0].text = text
                        else:
                            run = p.add_run()
                            run.text = text
                    else:
                        # Add new paragraph — then copy template formatting
                        p = tf.add_paragraph()
                        # Use add_run instead of p.text to preserve pPr (Pitfall 5)
                        run = p.add_run()
                        run.text = text

                        # Copy paragraph properties from first paragraph
                        if first_pPr is not None:
                            copy_paragraph_properties(first_pPr, p._element)

                        # Copy run properties from template's first run
                        if first_run_elem is not None:
                            target_r = p._element.find(qn("a:r"))
                            if target_r is not None:
                                copy_run_properties(first_run_elem, target_r)

                    # Set paragraph level
                    p.level = item_level

                    # Apply per-item bullet override if specified
                    effective_bullet = item_bullet_spec or list_bullet_spec
                    if effective_bullet is not None:
                        apply_bullet_type(p._element, effective_bullet, first_pPr)

                    # Apply per-item formatting overrides
                    if item_overrides:
                        self._apply_run_overrides(p, item_overrides)

            return True

        except Exception as e:
            log.warning("Could not set content in placeholder %d: %s", idx, e)
            return False

    def _flatten_content_list(
        self, items: list, base_level: int = 0
    ) -> list[tuple[str, int, dict | None, BulletSpec | None]]:
        """Flatten a content list into (text, level, overrides, bullet_spec) tuples.

        Handles:
        - Plain strings: "item" → (text, base_level, None, None)
        - Dicts: {"text": "item", "level": 1, "bold": True} → (text, level, overrides, bullet_spec)
        - Nested lists: ["parent", ["child1", "child2"]] → flattened with incremented levels

        Returns list of (text, level, overrides_dict_or_None, BulletSpec_or_None) tuples.
        """
        result = []
        for item in items:
            if isinstance(item, list):
                # Nested list = sub-bullets at next level
                result.extend(self._flatten_content_list(item, base_level + 1))
            elif isinstance(item, dict):
                text = str(item.get("text", ""))
                item_level = item.get("level", base_level)

                # Extract formatting overrides
                override_keys = {"bold", "italic", "font_size", "color", "font_name", "underline"}
                overrides = {k: v for k, v in item.items() if k in override_keys}

                # Extract per-item bullet spec
                item_bullet = item.get("bullet")
                item_bullet_spec = None
                if item_bullet is not None:
                    if isinstance(item_bullet, str):
                        if len(item_bullet) == 1 and item_bullet not in ("a",):
                            item_bullet_spec = BulletSpec(type="bullet", char=item_bullet)
                        else:
                            item_bullet_spec = BulletSpec(type=item_bullet)

                result.append((text, item_level, overrides if overrides else None, item_bullet_spec))
            else:
                result.append((str(item), base_level, None, None))
        return result

    def _apply_run_overrides(self, paragraph, overrides: dict) -> None:
        """Apply formatting overrides to all runs in a paragraph.

        Supports: bold, italic, underline, font_size (points), font_name, color (hex or theme name).
        """
        if not paragraph.runs:
            return

        for run in paragraph.runs:
            if "bold" in overrides:
                run.font.bold = overrides["bold"]
            if "italic" in overrides:
                run.font.italic = overrides["italic"]
            if "underline" in overrides:
                run.font.underline = overrides["underline"]
            if "font_size" in overrides:
                from pptx.util import Pt

                run.font.size = Pt(overrides["font_size"])
            if "font_name" in overrides:
                run.font.name = overrides["font_name"]
            if "color" in overrides:
                color_val = overrides["color"]
                if isinstance(color_val, str):
                    if color_val.startswith("#"):
                        # Hex color
                        from pptx.dml.color import RGBColor

                        hex_str = color_val.lstrip("#")
                        run.font.color.rgb = RGBColor(
                            int(hex_str[0:2], 16),
                            int(hex_str[2:4], 16),
                            int(hex_str[4:6], 16),
                        )
                    else:
                        # Theme color name — resolve via ThemeColors if available
                        # For now, log and skip (theme color on runs requires MSO_THEME_COLOR)
                        log.debug("Theme color '%s' on run override — skipping (use hex)", color_val)

    @staticmethod
    def _enable_autofit(placeholder) -> None:
        """Enable normAutofit (shrink-to-fit) on a placeholder's bodyPr."""
        txBody = placeholder._element.find(qn("p:txBody"))
        if txBody is None:
            return

        bodyPr = txBody.find(qn("a:bodyPr"))
        if bodyPr is None:
            return

        # Remove any existing autofit settings
        for child in list(bodyPr):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag in ("noAutofit", "normAutofit", "spAutoFit"):
                bodyPr.remove(child)

        # Add normAutofit (shrink text to fit)
        from lxml import etree

        etree.SubElement(bodyPr, qn("a:normAutofit"))


class SlideBuilder:
    """
    Build slides from a template, providing SlideProxy for each added slide.
    """

    def __init__(self, presentation: Presentation, template_info: TemplateInfo):
        self.prs = presentation
        self.template = template_info

    def add_slide(self, layout_name: str) -> SlideProxy:
        """Add a new slide using the named layout and return a SlideProxy."""
        layout_info = self.template.get_layout(layout_name)
        if layout_info is None:
            layout_info = self.template.find_layout(layout_name)
        if layout_info is None:
            raise ValueError(f"Layout '{layout_name}' not found in template '{self.template.name}'")

        # Get the actual pptx layout object
        pptx_layout = self.prs.slide_layouts[layout_info.index]
        slide = self.prs.slides.add_slide(pptx_layout)

        return SlideProxy(slide, layout_info)

    def save(self, path: str | Path) -> Path:
        """Save the presentation to disk."""
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        return path
