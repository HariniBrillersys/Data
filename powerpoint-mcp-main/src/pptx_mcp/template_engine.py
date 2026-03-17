"""
Template Engine - Load and introspect PowerPoint templates (.potx/.pptx).

Handles .potx conversion, caches loaded templates, and provides deep
introspection of every layout's placeholders with their exact indices,
types, positions and sizes.
"""

from __future__ import annotations

import hashlib
import json
import logging
import tempfile
import time
import zipfile
from dataclasses import asdict, dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.oxml.ns import qn

from .theme_colors import ThemeColors

log = logging.getLogger(__name__)

_CACHE_FILENAME = ".template_cache.json"


def resolve_font_size_from_xml(placeholder, default_pt: float = 14.0, _depth: int = 0) -> float:
    """
    Walk the XML inheritance chain to find the effective font size (in pt).

    Resolution order:
      1. lstStyle on the placeholder's txBody (layout-defined, most reliable)
      2. Paragraph-level defRPr
      3. Layout placeholder's lstStyle (from the slide's layout — only for
         slide-level placeholders, skipped for layout placeholders to avoid
         infinite recursion)
      4. Type-based defaults (title=20.0, subtitle=16.0, body=14.0)

    Args:
        placeholder: python-pptx placeholder object
        default_pt: Fallback default if no font size found (default: 14.0)
        _depth: Internal recursion guard (do not set manually)

    Returns:
        Font size in points
    """
    elem = placeholder._element

    # 1. Check lstStyle in the placeholder's own text body
    for ns_prefix in ("p:txBody", "a:txBody"):
        txBody = elem.find(qn(ns_prefix))
        if txBody is not None:
            lstStyle = txBody.find(qn("a:lstStyle"))
            if lstStyle is not None:
                for lvl in lstStyle:
                    defRPr = lvl.find(qn("a:defRPr"))
                    if defRPr is not None:
                        sz = defRPr.get("sz")
                        if sz:
                            return int(sz) / 100

    # 2. Check paragraph-level defRPr
    for pPr in elem.iter(qn("a:pPr")):
        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is not None:
            sz = defRPr.get("sz")
            if sz:
                return int(sz) / 100

    # 3. Try the layout's corresponding placeholder (slide-level only).
    #    Skip when already on a layout placeholder (part is SlideLayoutPart)
    #    to avoid infinite self-recursion, and cap depth as a safety net.
    if _depth < 2:
        try:
            from pptx.parts.slidelayout import SlideLayoutPart

            slide_part = placeholder.part
            if hasattr(slide_part, "slide_layout") and not isinstance(slide_part, SlideLayoutPart):
                layout = slide_part.slide_layout
                ph_idx = placeholder.placeholder_format.idx
                for layout_ph in layout.placeholders:
                    if layout_ph.placeholder_format.idx == ph_idx:
                        return resolve_font_size_from_xml(layout_ph, default_pt, _depth + 1)
        except Exception:
            pass

    # 4. Fallback defaults by placeholder type
    try:
        ph_type = str(placeholder.placeholder_format.type)
        if "TITLE" in ph_type:
            return 20.0
        elif "SUBTITLE" in ph_type:
            return 16.0
    except Exception:
        pass

    return default_pt


class PlaceholderRole(str, Enum):
    """Semantic role of a placeholder within a layout."""

    TITLE = "title"
    SUBTITLE = "subtitle"
    BODY = "body"
    CONTENT = "content"
    CONTENT_LEFT = "content_left"
    CONTENT_RIGHT = "content_right"
    CONTENT_1 = "content_1"
    CONTENT_2 = "content_2"
    CONTENT_3 = "content_3"
    PICTURE = "picture"
    SLIDE_NUMBER = "slide_number"
    FOOTER = "footer"
    DATE = "date"
    NOTICE = "notice"
    UNKNOWN = "unknown"


@dataclass
class PlaceholderInfo:
    """Detailed information about a single placeholder in a layout."""

    idx: int
    name: str
    type_str: str
    role: PlaceholderRole
    left: int  # EMU
    top: int  # EMU
    width: int  # EMU
    height: int  # EMU
    hint_text: str = ""
    default_font_size: float = 0.0
    font_family: str = ""
    text_alignment: str = "left"
    color_scheme: dict[str, str] = field(default_factory=dict)
    text_styles: dict[str, Any] = field(default_factory=dict)
    is_primary: bool = False
    semantic_role_hint: str = ""
    formatting_recommendation: str = "preserve"
    has_crop_geometry: bool = False
    max_comfortable_words: int = 0
    max_comfortable_lines: int = 0
    max_bullet_items: int = 0

    @property
    def area(self) -> int:
        return self.width * self.height

    @property
    def visual_priority(self) -> int:
        """Numeric priority for visual hierarchy (higher = more prominent)."""
        # Primary placeholders get highest priority
        if self.is_primary:
            return 100
        # Title-like roles get high priority
        if self.role in (PlaceholderRole.TITLE, PlaceholderRole.SUBTITLE):
            return 80
        # Body text gets medium priority
        if self.role == PlaceholderRole.BODY:
            return 60
        # Content areas scaled by size
        return max(40, min(60, int(self.area / 1_000_000)))


@dataclass
class LayoutInfo:
    """Detailed information about a slide layout."""

    name: str
    index: int
    placeholders: list[PlaceholderInfo] = field(default_factory=list)

    @property
    def has_title(self) -> bool:
        return any(p.role == PlaceholderRole.TITLE for p in self.placeholders)

    @property
    def has_content(self) -> bool:
        return any(p.role == PlaceholderRole.CONTENT for p in self.placeholders)

    @property
    def has_picture(self) -> bool:
        return any(p.role == PlaceholderRole.PICTURE for p in self.placeholders)

    @property
    def content_count(self) -> int:
        return sum(
            1
            for p in self.placeholders
            if p.role
            in (
                PlaceholderRole.CONTENT,
                PlaceholderRole.CONTENT_LEFT,
                PlaceholderRole.CONTENT_RIGHT,
                PlaceholderRole.CONTENT_1,
                PlaceholderRole.CONTENT_2,
                PlaceholderRole.CONTENT_3,
            )
        )

    @property
    def design_intent(self) -> str:
        """
        Classify layout by design intent based on placeholder composition.

        Returns one of: title_heavy, content_heavy, balanced, visual_heavy, divider
        """
        title_area = sum(p.area for p in self.placeholders if p.role == PlaceholderRole.TITLE)
        content_area = sum(
            p.area
            for p in self.placeholders
            if p.role
            in (
                PlaceholderRole.CONTENT,
                PlaceholderRole.CONTENT_LEFT,
                PlaceholderRole.CONTENT_RIGHT,
                PlaceholderRole.BODY,
            )
        )
        picture_area = sum(p.area for p in self.placeholders if p.role == PlaceholderRole.PICTURE)

        # Estimate total slide area from placeholder composition
        total_area = max(p.area for p in self.placeholders) * 3 if self.placeholders else 1

        title_pct = title_area / total_area if total_area > 0 else 0
        content_pct = content_area / total_area if total_area > 0 else 0
        picture_pct = picture_area / total_area if total_area > 0 else 0

        # Classification thresholds
        if title_pct > 0.5:
            return "divider"  # Title dominates (section dividers)
        elif picture_pct > 0.4:
            return "visual_heavy"  # Heavy on images
        elif content_pct > 0.5:
            return "content_heavy"  # Heavy on text content
        elif title_pct > 0.15 and content_pct > 0.15:
            return "balanced"  # Both title and content present
        elif title_pct > 0.1:
            return "title_heavy"  # Title present with minimal content
        else:
            return "balanced"  # Default classification

    def get_by_role(self, role: PlaceholderRole) -> Optional[PlaceholderInfo]:
        for p in self.placeholders:
            if p.role == role:
                return p
        return None

    def get_all_by_role(self, role: PlaceholderRole) -> list[PlaceholderInfo]:
        return [p for p in self.placeholders if p.role == role]

    def get_fillable_placeholders(self) -> list[PlaceholderInfo]:
        skip = {PlaceholderRole.SLIDE_NUMBER, PlaceholderRole.FOOTER, PlaceholderRole.DATE}
        return [p for p in self.placeholders if p.role not in skip]


@dataclass
class TemplateInfo:
    """Complete information about a loaded template."""

    name: str
    path: Path
    layouts: list[LayoutInfo] = field(default_factory=list)
    slide_width: int = 0
    slide_height: int = 0
    author: str = ""
    title: str = ""
    theme_colors: ThemeColors = field(default_factory=ThemeColors.uptimize_defaults)

    def get_layout(self, name: str) -> Optional[LayoutInfo]:
        for layout in self.layouts:
            if layout.name == name:
                return layout
        return None

    def find_layout(self, keyword: str) -> Optional[LayoutInfo]:
        keyword_lower = keyword.lower()
        for layout in self.layouts:
            if keyword_lower in layout.name.lower():
                return layout
        return None

    def layout_names(self) -> list[str]:
        return [layout.name for layout in self.layouts]

    def to_cache_dict(self) -> dict:
        """Serialize to a JSON-compatible dict for disk caching."""
        return {
            "name": self.name,
            "path": str(self.path),
            "slide_width": self.slide_width,
            "slide_height": self.slide_height,
            "author": self.author,
            "title": self.title,
            "theme_colors": asdict(self.theme_colors),
            "layouts": [
                {
                    "name": layout.name,
                    "index": layout.index,
                    "placeholders": [
                        {
                            **{k: v for k, v in asdict(ph).items() if k != "role"},
                            "role": ph.role.value,
                        }
                        for ph in layout.placeholders
                    ],
                }
                for layout in self.layouts
            ],
        }

    @classmethod
    def from_cache_dict(cls, data: dict) -> TemplateInfo:
        """Deserialize from a cached JSON dict."""
        theme_colors = ThemeColors(**data["theme_colors"])
        layouts = []
        for ld in data["layouts"]:
            phs = []
            for pd in ld["placeholders"]:
                pd["role"] = PlaceholderRole(pd["role"])
                phs.append(PlaceholderInfo(**pd))
            layouts.append(LayoutInfo(name=ld["name"], index=ld["index"], placeholders=phs))
        return cls(
            name=data["name"],
            path=Path(data["path"]),
            layouts=layouts,
            slide_width=data["slide_width"],
            slide_height=data["slide_height"],
            author=data.get("author", ""),
            title=data.get("title", ""),
            theme_colors=theme_colors,
        )


class TemplateEngine:
    """
    Load, cache, and introspect PowerPoint templates.

    Two-tier design:
      - list_available() returns template filenames instantly (no parsing)
      - get_template() deep-analyzes on first access, with JSON disk cache

    Core design: we NEVER create new textboxes. All content goes into
    existing placeholders identified by their idx.
    """

    def __init__(self, templates_dir: str | Path = "templates"):
        self.templates_dir = Path(templates_dir)
        self._cache: dict[str, TemplateInfo] = {}
        self._temp_files: list[Path] = []
        self._disk_cache: dict[str, dict] = {}
        self._disk_cache_path = self.templates_dir / _CACHE_FILENAME
        self._load_disk_cache()

    def _load_disk_cache(self) -> None:
        """Load the JSON disk cache if it exists."""
        if self._disk_cache_path.exists():
            try:
                with open(self._disk_cache_path) as f:
                    self._disk_cache = json.load(f)
                log.debug("Loaded template cache from %s", self._disk_cache_path)
            except Exception as e:
                log.debug("Could not load template cache: %s", e)
                self._disk_cache = {}

    def _save_disk_cache(self) -> None:
        """Persist the JSON disk cache to disk."""
        try:
            with open(self._disk_cache_path, "w") as f:
                json.dump(self._disk_cache, f, indent=2)
            log.debug("Saved template cache to %s", self._disk_cache_path)
        except Exception as e:
            log.warning("Could not save template cache: %s", e)

    @staticmethod
    def _file_hash(path: Path) -> str:
        """Compute a fast MD5 hash of a file for cache invalidation."""
        h = hashlib.md5()
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(1 << 20), b""):  # 1MB chunks
                h.update(chunk)
        return h.hexdigest()

    def _is_cache_valid(self, path: Path) -> bool:
        """Check if disk cache entry is still valid for a template file."""
        stem = path.stem
        if stem not in self._disk_cache:
            return False
        entry = self._disk_cache[stem]
        return entry.get("hash") == self._file_hash(path)

    def list_available(self) -> list[dict]:
        """List template files instantly without deep-scanning.

        Returns a list of dicts with name, file, and whether it's been analyzed.
        This is the fast path for 'which templates do you have?'
        """
        if not self.templates_dir.exists():
            return []

        templates = []
        for ext in ("*.potx", "*.pptx"):
            for path in sorted(self.templates_dir.glob(ext)):
                templates.append(
                    {
                        "name": path.stem,
                        "file": path.name,
                        "analyzed": path.stem in self._cache,
                    }
                )
        return templates

    def scan(self) -> list[str]:
        """Scan templates directory and register all .potx and .pptx templates.

        Uses disk cache to avoid re-analyzing unchanged templates.
        """
        if not self.templates_dir.exists():
            return []

        names = []
        cache_dirty = False
        for ext in ("*.potx", "*.pptx"):
            for path in sorted(self.templates_dir.glob(ext)):
                if path.stem in self._cache:
                    names.append(path.stem)
                    continue

                # Try disk cache first
                if self._is_cache_valid(path):
                    try:
                        info = TemplateInfo.from_cache_dict(self._disk_cache[path.stem]["data"])
                        self._cache[info.name] = info
                        names.append(info.name)
                        log.debug("Loaded '%s' from disk cache", info.name)
                        continue
                    except Exception as e:
                        log.debug("Cache deserialization failed for '%s': %s", path.stem, e)

                # Full analysis needed
                start = time.monotonic()
                info = self._analyze_template(path)
                elapsed = time.monotonic() - start
                self._cache[info.name] = info
                names.append(info.name)
                log.info("Analyzed '%s' in %.2fs", info.name, elapsed)

                # Update disk cache
                self._disk_cache[path.stem] = {"hash": self._file_hash(path), "data": info.to_cache_dict()}
                cache_dirty = True

        if cache_dirty:
            self._save_disk_cache()

        return names

    def ensure_loaded(self, name: str) -> Optional[TemplateInfo]:
        """Ensure a specific template is deep-analyzed and return it.

        Loads from disk cache if available, otherwise analyzes and caches.
        """
        if name in self._cache:
            return self._cache[name]

        # Find the file
        path = self._find_template_file(name)
        if path is None:
            return None

        # Try disk cache
        if self._is_cache_valid(path):
            try:
                info = TemplateInfo.from_cache_dict(self._disk_cache[path.stem]["data"])
                self._cache[info.name] = info
                log.debug("Loaded '%s' from disk cache", info.name)
                return info
            except Exception as e:
                log.debug("Cache deserialization failed for '%s': %s", path.stem, e)

        # Full analysis
        start = time.monotonic()
        info = self._analyze_template(path)
        elapsed = time.monotonic() - start
        self._cache[info.name] = info
        log.info("Analyzed '%s' in %.2fs", info.name, elapsed)

        # Update disk cache
        self._disk_cache[path.stem] = {"hash": self._file_hash(path), "data": info.to_cache_dict()}
        self._save_disk_cache()

        return info

    def _find_template_file(self, name: str) -> Optional[Path]:
        """Find a template file by stem name."""
        for ext in (".potx", ".pptx"):
            path = self.templates_dir / f"{name}{ext}"
            if path.exists():
                return path
        return None

    def register_template(self, path: str | Path, name: Optional[str] = None) -> TemplateInfo:
        """Register a single template file at runtime."""
        path = Path(path)
        if not path.exists():
            raise FileNotFoundError(f"Template file not found: {path}")
        if path.suffix not in (".potx", ".pptx"):
            raise ValueError(f"Template must be .potx or .pptx, got: {path.suffix}")

        info = self._analyze_template(path)
        if name:
            info.name = name
        self._cache[info.name] = info

        # Update disk cache
        self._disk_cache[path.stem] = {"hash": self._file_hash(path), "data": info.to_cache_dict()}
        self._save_disk_cache()

        return info

    def get_template(self, name: str) -> Optional[TemplateInfo]:
        """Get a template by name, analyzing on demand if needed."""
        return self.ensure_loaded(name)

    def list_templates(self) -> list[str]:
        """Return names of all analyzed (in-memory) templates."""
        return list(self._cache.keys())

    def open_presentation(self, template_name: str) -> Presentation:
        """Open a fresh Presentation from a template, ready for slide creation."""
        info = self._cache.get(template_name)
        if not info:
            raise ValueError(f"Template '{template_name}' not found. Available: {self.list_templates()}")

        path = info.path
        if path.suffix == ".potx":
            pptx_path = self._convert_potx(path)
            prs = Presentation(str(pptx_path))
        else:
            prs = Presentation(str(path))

        # Remove any example slides that come with the template
        while len(prs.slides) > 0:
            r_id = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[0]

        return prs

    def _convert_potx(self, potx_path: Path) -> Path:
        """Convert .potx to .pptx by rewriting the content type in the ZIP."""
        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        tmp_path = Path(tmp.name)
        tmp.close()

        with zipfile.ZipFile(potx_path, "r") as zin:
            with zipfile.ZipFile(tmp_path, "w") as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    if item.filename == "[Content_Types].xml":
                        data = data.replace(
                            b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                            b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                        )
                    zout.writestr(item, data)

        self._temp_files.append(tmp_path)
        return tmp_path

    def _analyze_template(self, path: Path) -> TemplateInfo:
        """Deep-analyze a template file and extract all layout/placeholder info."""
        if path.suffix == ".potx":
            pptx_path = self._convert_potx(path)
            prs = Presentation(str(pptx_path))
        else:
            prs = Presentation(str(path))

        # Extract theme colors
        theme_colors = ThemeColors.from_template(path)

        info = TemplateInfo(
            name=path.stem,
            path=path,
            slide_width=prs.slide_width,
            slide_height=prs.slide_height,
            theme_colors=theme_colors,
        )

        try:
            info.author = prs.core_properties.author or ""
            info.title = prs.core_properties.title or ""
        except Exception:
            pass

        for layout_idx, layout in enumerate(prs.slide_layouts):
            layout_info = self._analyze_layout(layout, layout_idx)
            info.layouts.append(layout_info)

        return info

    def _analyze_layout(self, layout, layout_idx: int) -> LayoutInfo:
        """Analyze a single layout and classify all its placeholders."""
        layout_info = LayoutInfo(name=layout.name, index=layout_idx)

        raw_phs = []
        ph_objects = {}  # Store placeholder objects by idx for extraction
        for ph in layout.placeholders:
            try:
                idx = ph.placeholder_format.idx
                raw_phs.append(
                    {
                        "idx": idx,
                        "name": ph.name,
                        "type_str": str(ph.placeholder_format.type),
                        "left": ph.left or 0,
                        "top": ph.top or 0,
                        "width": ph.width or 0,
                        "height": ph.height or 0,
                    }
                )
                ph_objects[idx] = ph  # Store for extraction
            except Exception:
                continue

        object_phs = []
        for raw in raw_phs:
            role = self._classify_placeholder(raw)
            ph_obj = ph_objects.get(raw["idx"])
            if role == PlaceholderRole.CONTENT:
                object_phs.append(raw)
            else:
                font_size = self._extract_font_size(ph_obj) if ph_obj else 0.0
                max_words, max_lines, max_bullets = self._calculate_capacity(
                    raw["width"], raw["height"], font_size or 14.0, role
                )
                layout_info.placeholders.append(
                    PlaceholderInfo(
                        idx=raw["idx"],
                        name=raw["name"],
                        type_str=raw["type_str"],
                        role=role,
                        left=raw["left"],
                        top=raw["top"],
                        width=raw["width"],
                        height=raw["height"],
                        hint_text=self._extract_hint_text(ph_obj) if ph_obj else "",
                        default_font_size=font_size,
                        font_family=self._extract_font_family(ph_obj) if ph_obj else "",
                        text_alignment=self._extract_alignment(ph_obj) if ph_obj else "left",
                        color_scheme=self._extract_color_scheme(ph_obj) if ph_obj else {},
                        text_styles=self._extract_text_styles(ph_obj) if ph_obj else {},
                        has_crop_geometry=(
                            self._detect_crop_geometry(ph_obj) if ph_obj and role == PlaceholderRole.PICTURE else False
                        ),
                        max_comfortable_words=max_words,
                        max_comfortable_lines=max_lines,
                        max_bullet_items=max_bullets,
                    )
                )

        # Handle OBJECT placeholders: assign spatial roles if multiple
        if len(object_phs) == 1:
            ph = object_phs[0]
            ph_obj = ph_objects.get(ph["idx"])
            font_size = self._extract_font_size(ph_obj) if ph_obj else 0.0
            max_words, max_lines, max_bullets = self._calculate_capacity(
                ph["width"], ph["height"], font_size or 14.0, PlaceholderRole.CONTENT
            )
            layout_info.placeholders.append(
                PlaceholderInfo(
                    idx=ph["idx"],
                    name=ph["name"],
                    type_str=ph["type_str"],
                    role=PlaceholderRole.CONTENT,
                    left=ph["left"],
                    top=ph["top"],
                    width=ph["width"],
                    height=ph["height"],
                    hint_text=self._extract_hint_text(ph_obj) if ph_obj else "",
                    default_font_size=font_size,
                    font_family=self._extract_font_family(ph_obj) if ph_obj else "",
                    text_alignment=self._extract_alignment(ph_obj) if ph_obj else "left",
                    color_scheme=self._extract_color_scheme(ph_obj) if ph_obj else {},
                    text_styles=self._extract_text_styles(ph_obj) if ph_obj else {},
                    has_crop_geometry=False,
                    max_comfortable_words=max_words,
                    max_comfortable_lines=max_lines,
                    max_bullet_items=max_bullets,
                )
            )
        elif len(object_phs) == 2:
            sorted_phs = sorted(object_phs, key=lambda p: p["left"])
            area_0 = sorted_phs[0]["width"] * sorted_phs[0]["height"]
            area_1 = sorted_phs[1]["width"] * sorted_phs[1]["height"]

            if area_1 < area_0 * 0.5:
                roles = [PlaceholderRole.CONTENT, PlaceholderRole.NOTICE]
            elif area_0 < area_1 * 0.5:
                roles = [PlaceholderRole.NOTICE, PlaceholderRole.CONTENT]
            else:
                roles = [PlaceholderRole.CONTENT_LEFT, PlaceholderRole.CONTENT_RIGHT]

            for ph, role in zip(sorted_phs, roles):
                ph_obj = ph_objects.get(ph["idx"])
                font_size = self._extract_font_size(ph_obj) if ph_obj else 0.0
                max_words, max_lines, max_bullets = self._calculate_capacity(
                    ph["width"], ph["height"], font_size or 14.0, role
                )
                layout_info.placeholders.append(
                    PlaceholderInfo(
                        idx=ph["idx"],
                        name=ph["name"],
                        type_str=ph["type_str"],
                        role=role,
                        left=ph["left"],
                        top=ph["top"],
                        width=ph["width"],
                        height=ph["height"],
                        hint_text=self._extract_hint_text(ph_obj) if ph_obj else "",
                        default_font_size=font_size,
                        font_family=self._extract_font_family(ph_obj) if ph_obj else "",
                        text_alignment=self._extract_alignment(ph_obj) if ph_obj else "left",
                        color_scheme=self._extract_color_scheme(ph_obj) if ph_obj else {},
                        text_styles=self._extract_text_styles(ph_obj) if ph_obj else {},
                        has_crop_geometry=False,
                        max_comfortable_words=max_words,
                        max_comfortable_lines=max_lines,
                        max_bullet_items=max_bullets,
                    )
                )
        elif len(object_phs) >= 3:
            sorted_phs = sorted(object_phs, key=lambda p: p["left"])
            numbered_roles = [
                PlaceholderRole.CONTENT_1,
                PlaceholderRole.CONTENT_2,
                PlaceholderRole.CONTENT_3,
            ]
            for i, ph in enumerate(sorted_phs):
                role = numbered_roles[i] if i < len(numbered_roles) else PlaceholderRole.CONTENT
                ph_obj = ph_objects.get(ph["idx"])
                font_size = self._extract_font_size(ph_obj) if ph_obj else 0.0
                max_words, max_lines, max_bullets = self._calculate_capacity(
                    ph["width"], ph["height"], font_size or 14.0, role
                )
                layout_info.placeholders.append(
                    PlaceholderInfo(
                        idx=ph["idx"],
                        name=ph["name"],
                        type_str=ph["type_str"],
                        role=role,
                        left=ph["left"],
                        top=ph["top"],
                        width=ph["width"],
                        height=ph["height"],
                        hint_text=self._extract_hint_text(ph_obj) if ph_obj else "",
                        default_font_size=font_size,
                        font_family=self._extract_font_family(ph_obj) if ph_obj else "",
                        text_alignment=self._extract_alignment(ph_obj) if ph_obj else "left",
                        color_scheme=self._extract_color_scheme(ph_obj) if ph_obj else {},
                        text_styles=self._extract_text_styles(ph_obj) if ph_obj else {},
                        has_crop_geometry=False,
                        max_comfortable_words=max_words,
                        max_comfortable_lines=max_lines,
                        max_bullet_items=max_bullets,
                    )
                )

        # Populate semantic hints and formatting recommendations
        for ph in layout_info.placeholders:
            ph.semantic_role_hint = self._infer_semantic_role(ph.hint_text, ph.default_font_size, ph.role)
            ph.formatting_recommendation = self._recommend_formatting(ph.hint_text, ph.default_font_size, ph.role)

        # Determine visual hierarchy (sets is_primary flags)
        self._determine_visual_hierarchy(layout_info.placeholders)

        return layout_info

    def _classify_placeholder(self, raw: dict) -> PlaceholderRole:
        """Classify a placeholder by its type string into a semantic role."""
        t = raw["type_str"]

        if "CENTER_TITLE" in t or ("TITLE" in t and "SUBTITLE" not in t):
            return PlaceholderRole.TITLE
        if "SUBTITLE" in t:
            return PlaceholderRole.SUBTITLE
        if "PICTURE" in t:
            return PlaceholderRole.PICTURE
        if "SLIDE_NUMBER" in t:
            return PlaceholderRole.SLIDE_NUMBER
        if "FOOTER" in t:
            return PlaceholderRole.FOOTER
        if "DATE" in t:
            return PlaceholderRole.DATE

        if "OBJECT" in t:
            return PlaceholderRole.CONTENT
        if "BODY" in t:
            if raw["height"] > 2000000:
                return PlaceholderRole.CONTENT
            return PlaceholderRole.BODY

        return PlaceholderRole.UNKNOWN

    def _extract_hint_text(self, placeholder) -> str:
        """Extract hint text with defensive parsing for PowerPoint 2013/2016/365 XML variations."""
        try:
            ph_element = placeholder._element

            # Try p:ph element first (layout/master level)
            ph = ph_element.find(qn("p:ph"))
            if ph is not None:
                hint = ph.get("text")
                if hint:
                    return hint

            # Try first paragraph text (common for layout placeholders)
            txBody = ph_element.find(qn("p:txBody"))
            if txBody is not None:
                first_p = txBody.find(qn("a:p"))
                if first_p is not None:
                    text = "".join(first_p.itertext())
                    if text:
                        return text.strip()

            return ""
        except Exception:
            return ""

    def _extract_font_size(self, placeholder) -> float:
        """Extract font size with defensive parsing for PowerPoint 2013/2016/365 XML variations."""
        try:
            return resolve_font_size_from_xml(placeholder, default_pt=14.0)
        except Exception:
            return 14.0

    def _extract_font_family(self, placeholder) -> str:
        """Extract font family with defensive parsing for PowerPoint 2013/2016/365 XML variations."""
        try:
            elem = placeholder._element
            txBody = elem.find(qn("p:txBody"))
            if txBody is None:
                return ""

            lstStyle = txBody.find(qn("a:lstStyle"))
            if lstStyle is not None:
                for lvl in lstStyle:
                    defRPr = lvl.find(qn("a:defRPr"))
                    if defRPr is not None:
                        latin = defRPr.find(qn("a:latin"))
                        if latin is not None:
                            typeface = latin.get("typeface")
                            if typeface:
                                return typeface

            # Check run properties
            for rPr in elem.iter(qn("a:rPr")):
                latin = rPr.find(qn("a:latin"))
                if latin is not None:
                    typeface = latin.get("typeface")
                    if typeface:
                        return typeface

            return ""
        except Exception:
            return ""

    def _extract_alignment(self, placeholder) -> str:
        """Extract alignment with defensive parsing for PowerPoint 2013/2016/365 XML variations."""
        try:
            txBody = placeholder._element.find(qn("p:txBody"))
            if txBody is None:
                return "left"

            # Check first paragraph's alignment
            first_p = txBody.find(qn("a:p"))
            if first_p is None:
                return "left"

            pPr = first_p.find(qn("a:pPr"))
            if pPr is None:
                return "left"

            algn = pPr.get("algn")
            if algn:
                # Map XML values to readable strings
                mapping = {"l": "left", "ctr": "center", "r": "right"}
                return mapping.get(algn, "left")

            return "left"
        except Exception:
            return "left"

    def _extract_color_scheme(self, placeholder) -> dict[str, str]:
        """Extract color scheme from placeholder XML."""
        try:
            elem = placeholder._element
            txBody = elem.find(qn("p:txBody"))
            if txBody is None:
                return {}

            result = {}

            # Check lstStyle for theme color references
            lstStyle = txBody.find(qn("a:lstStyle"))
            if lstStyle is not None:
                for lvl in lstStyle:
                    defRPr = lvl.find(qn("a:defRPr"))
                    if defRPr is not None:
                        solidFill = defRPr.find(qn("a:solidFill"))
                        if solidFill is not None:
                            # Check for theme color reference
                            schemeClr = solidFill.find(qn("a:schemeClr"))
                            if schemeClr is not None:
                                val = schemeClr.get("val")
                                if val:
                                    result["scheme_color"] = val

                            # Check for explicit RGB color
                            srgbClr = solidFill.find(qn("a:srgbClr"))
                            if srgbClr is not None:
                                val = srgbClr.get("val")
                                if val:
                                    result["rgb"] = f"#{val.upper()}"

                            if result:
                                return result

            return {}
        except Exception:
            return {}

    def _extract_text_styles(self, placeholder) -> dict[str, Any]:
        """Extract multi-level text styles from placeholder XML."""
        try:
            elem = placeholder._element
            txBody = elem.find(qn("p:txBody"))
            if txBody is None:
                return {}

            lstStyle = txBody.find(qn("a:lstStyle"))
            if lstStyle is None:
                return {}

            styles = {}

            # Extract up to 9 levels (lvl0-lvl8)
            for i in range(9):
                lvl_elem = lstStyle.find(qn(f"a:lvl{i}pPr"))
                if lvl_elem is not None:
                    level_styles = {}

                    # Extract indentation
                    mar_l = lvl_elem.get("marL")
                    if mar_l is not None:
                        level_styles["indent"] = int(mar_l)

                    indent_attr = lvl_elem.get("indent")
                    if indent_attr is not None:
                        level_styles["bullet_indent"] = int(indent_attr)

                    # Extract line spacing
                    lnSpc = lvl_elem.find(qn("a:lnSpc"))
                    if lnSpc is not None:
                        spcPct = lnSpc.find(qn("a:spcPct"))
                        if spcPct is not None:
                            val = spcPct.get("val")
                            if val:
                                level_styles["line_spacing"] = int(val)

                    # Extract space before/after (as child elements, not attributes)
                    spc_bef_elem = lvl_elem.find(qn("a:spcBef"))
                    if spc_bef_elem is not None:
                        spc_pts = spc_bef_elem.find(qn("a:spcPts"))
                        if spc_pts is not None:
                            val = spc_pts.get("val")
                            if val:
                                level_styles["space_before"] = int(val)  # hundredths of a point
                        else:
                            spc_pct = spc_bef_elem.find(qn("a:spcPct"))
                            if spc_pct is not None:
                                val = spc_pct.get("val")
                                if val:
                                    level_styles["space_before_pct"] = int(val)  # hundredths of percent

                    spc_aft_elem = lvl_elem.find(qn("a:spcAft"))
                    if spc_aft_elem is not None:
                        spc_pts = spc_aft_elem.find(qn("a:spcPts"))
                        if spc_pts is not None:
                            val = spc_pts.get("val")
                            if val:
                                level_styles["space_after"] = int(val)  # hundredths of a point
                        else:
                            spc_pct = spc_aft_elem.find(qn("a:spcPct"))
                            if spc_pct is not None:
                                val = spc_pct.get("val")
                                if val:
                                    level_styles["space_after_pct"] = int(val)  # hundredths of percent

                    # Extract text run properties (bold/italic)
                    defRPr = lvl_elem.find(qn("a:defRPr"))
                    if defRPr is not None:
                        bold = defRPr.get("b")
                        if bold is not None:
                            level_styles["bold"] = bold == "1"

                        italic = defRPr.get("i")
                        if italic is not None:
                            level_styles["italic"] = italic == "1"

                        underline = defRPr.get("u")
                        if underline is not None:
                            level_styles["underline"] = underline != "none"

                    # Extract bullet type elements
                    buChar = lvl_elem.find(qn("a:buChar"))
                    if buChar is not None:
                        level_styles["bullet_type"] = "char"
                        level_styles["bullet_char"] = buChar.get("char", "•")

                    buAutoNum = lvl_elem.find(qn("a:buAutoNum"))
                    if buAutoNum is not None:
                        level_styles["bullet_type"] = "autonum"
                        level_styles["bullet_scheme"] = buAutoNum.get("type", "arabicPeriod")
                        start_at = buAutoNum.get("startAt")
                        if start_at is not None:
                            level_styles["bullet_start_at"] = int(start_at)

                    buNone = lvl_elem.find(qn("a:buNone"))
                    if buNone is not None:
                        level_styles["bullet_type"] = "none"

                    # Extract bullet font if present
                    buFont = lvl_elem.find(qn("a:buFont"))
                    if buFont is not None:
                        font_typeface = buFont.get("typeface")
                        if font_typeface:
                            level_styles["bullet_font"] = font_typeface

                    if level_styles:
                        styles[f"level_{i}"] = level_styles

            return styles
        except Exception:
            return {}

    def _determine_visual_hierarchy(self, placeholders: list[PlaceholderInfo]) -> None:
        """Determine visual hierarchy by setting is_primary flag on most prominent placeholder.

        Modifies placeholders in-place. Primary = largest by area among top 30% of slide.
        Only sets is_primary=True for ONE placeholder per layout.
        """
        if not placeholders:
            return

        # Find placeholders in top 30% of slide (lowest top values)
        slide_height = 6858000  # Standard 16:9 slide height in EMU
        top_threshold = slide_height * 0.3
        top_placeholders = [p for p in placeholders if p.top <= top_threshold]

        if not top_placeholders:
            # If no placeholders in top 30%, consider all
            top_placeholders = placeholders

        # Find largest by area
        largest = max(top_placeholders, key=lambda p: p.area)

        # If tie, prefer TITLE > CONTENT > others
        candidates = [p for p in top_placeholders if p.area == largest.area]
        if len(candidates) > 1:
            # Sort by role priority
            role_priority = {
                PlaceholderRole.TITLE: 3,
                PlaceholderRole.CONTENT: 2,
                PlaceholderRole.CONTENT_LEFT: 2,
                PlaceholderRole.CONTENT_RIGHT: 2,
            }
            candidates.sort(key=lambda p: role_priority.get(p.role, 1), reverse=True)
            largest = candidates[0]

        largest.is_primary = True

    def _infer_semantic_role(self, hint_text: str, default_font_size: float, role: PlaceholderRole) -> str:
        """Infer semantic role from hint text patterns and font size.

        Returns semantic hint based on hint text patterns and visual properties.
        Defaults to role.value if no pattern matches.
        """
        hint_lower = hint_text.lower()

        # Detect chapter numbers (00, 01, 02, etc.)
        if hint_text.strip() in ("00", "01", "02", "03", "04", "05", "06", "07", "08", "09"):
            return "chapter_number"
        if hint_text.strip().isdigit():
            return "chapter_number"

        # Detect chapter/section titles
        if "chapter" in hint_lower:
            return "chapter_title"
        if "section" in hint_lower:
            return "section_title"

        # Font size-based inference
        if default_font_size >= 48:
            return "emphasis"
        if default_font_size > 0 and default_font_size < 12:
            return "detail"

        # Default to role name
        return role.value

    def _recommend_formatting(self, hint_text: str, default_font_size: float, role: PlaceholderRole) -> str:
        """Recommend formatting strategy for this placeholder.

        Returns "preserve" (keep exact font size) or "shrink_if_needed" (allow auto-fit).
        """
        # Titles and subtitles should preserve size (shouldn't wrap)
        if role in (PlaceholderRole.TITLE, PlaceholderRole.SUBTITLE):
            return "preserve"

        # Decorative numbers should preserve size
        if hint_text.strip().isdigit() or hint_text.strip() in (
            "00",
            "01",
            "02",
            "03",
            "04",
            "05",
            "06",
            "07",
            "08",
            "09",
        ):
            return "preserve"

        # Large deliberate fonts should preserve
        if default_font_size >= 32:
            return "preserve"

        # Everything else can shrink if needed
        return "shrink_if_needed"

    def _calculate_capacity(
        self, width_emu: int, height_emu: int, font_size_pt: float, role: PlaceholderRole
    ) -> tuple[int, int, int]:
        """Calculate comfortable text capacity for a placeholder.

        Args:
            width_emu: Placeholder width in EMU
            height_emu: Placeholder height in EMU
            font_size_pt: Font size in points
            role: Placeholder role

        Returns:
            Tuple of (max_comfortable_words, max_comfortable_lines, max_bullet_items)
        """
        from .overflow import (
            _AVG_CHAR_WIDTH_EM,
            _DEFAULT_MARGIN_LR,
            _DEFAULT_MARGIN_TB,
            _EMU_PER_PT,
            _LINE_HEIGHT_FACTOR,
        )

        # Account for margins
        usable_width_pt = (width_emu - _DEFAULT_MARGIN_LR * 2) / _EMU_PER_PT
        usable_height_pt = (height_emu - _DEFAULT_MARGIN_TB * 2) / _EMU_PER_PT

        # Calculate line capacity
        char_width_pt = font_size_pt * _AVG_CHAR_WIDTH_EM
        chars_per_line = max(1, int(usable_width_pt / char_width_pt))

        line_height_pt = font_size_pt * _LINE_HEIGHT_FACTOR
        max_lines = max(1, int(usable_height_pt / line_height_pt))

        # Apply comfort factor
        max_comfortable_lines = int(max_lines * 0.85)

        # Calculate word capacity (avg 5 chars per word + space)
        avg_words_per_line = chars_per_line / 6
        max_comfortable_words = int(max_comfortable_lines * avg_words_per_line)

        # Calculate bullet capacity
        if role in (PlaceholderRole.TITLE, PlaceholderRole.SUBTITLE):
            max_bullet_items = 0  # Not intended for bullets
        else:
            # Cap at 7 for 6x6 guideline with slight flex
            max_bullet_items = min(max_comfortable_lines, 7)

        return max_comfortable_words, max_comfortable_lines, max_bullet_items

    def _detect_crop_geometry(self, placeholder) -> bool:
        """Detect if a placeholder has crop/mask geometry that would clip images.

        A placeholder has crop geometry if ANY of these are true:
        1. Custom geometry: shape has prstGeom with prst != "rect"
        2. Freeform shape: shape contains custGeom
        3. Pre-defined crop: blipFill contains srcRect with non-zero crop values
        4. Significant rounded corners: roundRect with large corner radii

        Args:
            placeholder: python-pptx placeholder object

        Returns:
            True if placeholder would crop/mask images
        """
        try:
            elem = placeholder._element

            # Check shape properties
            spPr = elem.find(qn("p:spPr"))
            if spPr is None:
                return False

            # 1. Check for custom geometry (non-rectangular shapes)
            prstGeom = spPr.find(qn("a:prstGeom"))
            if prstGeom is not None:
                prst = prstGeom.get("prst")
                if prst and prst != "rect":
                    # Any non-rectangular preset shape will crop
                    return True

            # 2. Check for freeform custom geometry
            custGeom = spPr.find(qn("a:custGeom"))
            if custGeom is not None:
                # Custom drawn shapes are masking paths
                return True

            # 3. Check blipFill for pre-defined crop
            blipFill = spPr.find(qn("p:blipFill"))
            if blipFill is not None:
                srcRect = blipFill.find(qn("a:srcRect"))
                if srcRect is not None:
                    # Check if any crop attribute is present and non-zero
                    for attr in ("l", "t", "r", "b"):
                        val = srcRect.get(attr)
                        if val and int(val) > 0:
                            return True

            # 4. Check for significant rounded corners
            # If prstGeom is roundRect, flag it (corner radius parsing is complex,
            # so we conservatively flag all roundRect shapes)
            if prstGeom is not None and prstGeom.get("prst") == "roundRect":
                return True

            return False

        except Exception as e:
            # Defensive: if we can't determine, assume no crop geometry
            log.debug("Could not detect crop geometry: %s", e)
            return False

    def cleanup(self):
        """Remove temporary files."""
        for path in self._temp_files:
            try:
                path.unlink()
            except OSError:
                pass
        self._temp_files.clear()

    def to_catalog(self) -> dict:
        """Export full template catalog as a serializable dict.

        DEPRECATED: This method returns minimal placeholder info (idx, role, name only).
        It does not include Phase 1 metadata fields (hint_text, font_size_pt, font_family,
        alignment, color_scheme, text_styles, visual_priority, is_primary, semantic_role_hint,
        formatting_recommendation, design_intent).

        Use get_template_layouts MCP tool instead for complete placeholder information.
        """
        catalog = {}
        for name, info in self._cache.items():
            catalog[name] = {
                "name": name,
                "author": info.author,
                "layouts": [],
            }
            for layout in info.layouts:
                layout_data = {
                    "name": layout.name,
                    "placeholders": [
                        {"idx": p.idx, "role": p.role.value, "name": p.name} for p in layout.get_fillable_placeholders()
                    ],
                }
                catalog[name]["layouts"].append(layout_data)
        return catalog
