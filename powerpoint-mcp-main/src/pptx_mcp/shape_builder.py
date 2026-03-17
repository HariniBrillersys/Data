"""
Shape Builder - Add decorative shapes on top of template placeholders.

Shapes enhance visual communication without replacing placeholder content.
All shapes use the template's theme colors for brand consistency.

Supported shape types:
  - arrow: Directional arrows (right, left, up, down)
  - callout: Text callout boxes
  - badge: Small label tags (e.g. "NEW", "IMPORTANT")
  - highlight: Semi-transparent overlay rectangles
  - connector: Lines between points
  - process_arrow: Chevron process flow arrows

Positioning is dynamic -- shapes are placed relative to slide dimensions
using named positions (top-right, bottom-left, center, etc.) or auto mode.
"""

from __future__ import annotations

import logging
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .theme_colors import ThemeColors

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Position resolver
# ---------------------------------------------------------------------------

# Named positions as fractions of (slide_width, slide_height)
_NAMED_POSITIONS = {
    "top-left": (0.05, 0.08),
    "top-center": (0.35, 0.08),
    "top-right": (0.75, 0.08),
    "center-left": (0.05, 0.42),
    "center": (0.35, 0.42),
    "center-right": (0.75, 0.42),
    "bottom-left": (0.05, 0.82),
    "bottom-center": (0.35, 0.82),
    "bottom-right": (0.75, 0.82),
    "below-title": (0.05, 0.22),
    "above-footer": (0.05, 0.88),
}


def _resolve_position(
    position: str | dict,
    slide_width: int,
    slide_height: int,
    shape_width: int,
    shape_height: int,
) -> tuple[int, int]:
    """Resolve a position spec to (left, top) in EMU.

    Args:
        position: Named position string or dict with 'left'/'top' in EMU.
        slide_width, slide_height: Slide dimensions in EMU.
        shape_width, shape_height: Shape dimensions in EMU.

    Returns:
        Tuple of (left, top) in EMU.
    """
    if isinstance(position, dict):
        return int(position.get("left", 0)), int(position.get("top", 0))

    pos = position.lower().strip()
    if pos in _NAMED_POSITIONS:
        frac_x, frac_y = _NAMED_POSITIONS[pos]
        left = int(slide_width * frac_x)
        top = int(slide_height * frac_y)
        return left, top

    # "auto" or unknown -- place in bottom-right area
    left = slide_width - shape_width - int(slide_width * 0.05)
    top = slide_height - shape_height - int(slide_height * 0.12)
    return left, top


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _choose_text_color(fill_hex: str) -> RGBColor:
    """Return white or dark text based on fill color luminance.

    Uses relative luminance formula: 0.299*R + 0.587*G + 0.114*B
    Dark text (#1D1D1D) for light backgrounds (luminance > 0.5)
    White text for dark backgrounds (luminance <= 0.5)
    """
    h = fill_hex.lstrip("#")
    r = int(h[0:2], 16)
    g = int(h[2:4], 16)
    b = int(h[4:6], 16)
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255

    if luminance > 0.5:
        return RGBColor(0x1D, 0x1D, 0x1D)  # Dark text for light backgrounds
    return RGBColor(0xFF, 0xFF, 0xFF)  # White text for dark backgrounds


# ---------------------------------------------------------------------------
# Individual shape builders
# ---------------------------------------------------------------------------


def _add_arrow(
    slide,
    direction: str = "right",
    label: str = "",
    color_hex: str = "#0F69AF",
    position: str | dict = "auto",
    slide_width: int = 0,
    slide_height: int = 0,
) -> dict:
    """Add a directional arrow shape."""
    shape_map = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    mso_shape = shape_map.get(direction.lower(), MSO_SHAPE.RIGHT_ARROW)

    # Size: arrows are typically wide and short (or tall and thin for vertical)
    if direction.lower() in ("up", "down"):
        width = Inches(0.8)
        height = Inches(1.5)
    else:
        width = Inches(1.8)
        height = Inches(0.7)

    left, top = _resolve_position(position, slide_width, slide_height, width, height)

    shape = slide.shapes.add_shape(mso_shape, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(color_hex)
    shape.line.fill.background()  # No border

    if label:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = _choose_text_color(color_hex)
        p.alignment = PP_ALIGN.CENTER

    return {"type": "arrow", "direction": direction, "label": label}


def _add_callout(
    slide,
    text: str = "",
    color_hex: str = "#0F69AF",
    position: str | dict = "auto",
    slide_width: int = 0,
    slide_height: int = 0,
) -> dict:
    """Add a text callout box."""
    width = Inches(2.5)
    height = Inches(1.0)

    left, top = _resolve_position(position, slide_width, slide_height, width, height)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(color_hex)
    shape.line.fill.background()

    # Adjust corner radius
    shape.adjustments[0] = 0.15

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = _choose_text_color(color_hex)
        p.alignment = PP_ALIGN.CENTER

    return {"type": "callout", "text": text}


def _add_badge(
    slide,
    text: str = "",
    color_hex: str = "#96D7D2",
    position: str | dict = "top-right",
    slide_width: int = 0,
    slide_height: int = 0,
) -> dict:
    """Add a small label/badge tag."""
    # Badge size: compact
    char_count = max(len(text), 3)
    width = Inches(max(0.8, min(2.0, char_count * 0.13)))
    height = Inches(0.35)

    left, top = _resolve_position(position, slide_width, slide_height, width, height)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(color_hex)
    shape.line.fill.background()

    shape.adjustments[0] = 0.35  # More rounded for badge feel

    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text.upper()
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = _choose_text_color(color_hex)
        p.alignment = PP_ALIGN.CENTER

    return {"type": "badge", "text": text}


def _add_highlight(
    slide,
    color_hex: str = "#96D7D2",
    opacity: float = 0.3,
    position: str | dict = "center",
    width_inches: float = 3.0,
    height_inches: float = 1.5,
    slide_width: int = 0,
    slide_height: int = 0,
) -> dict:
    """Add a semi-transparent highlight rectangle."""
    width = Inches(width_inches)
    height = Inches(height_inches)

    left, top = _resolve_position(position, slide_width, slide_height, width, height)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(color_hex)
    shape.line.fill.background()

    # Set transparency via XML (python-pptx doesn't expose this directly)
    from pptx.oxml.ns import qn

    solidFill = shape.fill._fill
    srgbClr = solidFill.find(qn("a:srgbClr"))
    if srgbClr is None:
        srgbClr = solidFill.find(f".//{qn('a:srgbClr')}")
    if srgbClr is not None:
        from lxml import etree

        alpha = etree.SubElement(srgbClr, qn("a:alpha"))
        alpha.set("val", str(int((1.0 - opacity) * 100000)))

    return {"type": "highlight", "opacity": opacity}


def _add_connector(
    slide,
    start: str | dict = "center-left",
    end: str | dict = "center-right",
    style: str = "solid",
    color_hex: str = "#0F69AF",
    weight_pt: float = 2.0,
    slide_width: int = 0,
    slide_height: int = 0,
) -> dict:
    """Add a connector line between two points."""
    start_x, start_y = _resolve_position(start, slide_width, slide_height, 0, 0)
    end_x, end_y = _resolve_position(end, slide_width, slide_height, 0, 0)

    from lxml import etree
    from pptx.oxml.ns import qn

    # Use the low-level API to add a line connector
    sp_tree = slide.shapes._spTree
    sp = etree.SubElement(sp_tree, qn("p:cxnSp"))

    # Non-visual properties
    nvCxnSpPr = etree.SubElement(sp, qn("p:nvCxnSpPr"))
    cNvPr = etree.SubElement(nvCxnSpPr, qn("p:cNvPr"))
    cNvPr.set("id", str(len(slide.shapes) + 100))
    cNvPr.set("name", "Connector")
    etree.SubElement(nvCxnSpPr, qn("p:cNvCxnSpPr"))
    etree.SubElement(nvCxnSpPr, qn("p:nvPr"))

    # Shape properties
    spPr = etree.SubElement(sp, qn("p:spPr"))

    # Position via xfrm
    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", str(min(start_x, end_x)))
    off.set("y", str(min(start_y, end_y)))
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(abs(end_x - start_x) or 1))
    ext.set("cy", str(abs(end_y - start_y) or 1))

    # If end is to the left or above start, flip
    if end_x < start_x:
        xfrm.set("flipH", "1")
    if end_y < start_y:
        xfrm.set("flipV", "1")

    # Preset geometry: straight connector
    prstGeom = etree.SubElement(spPr, qn("a:prstGeom"))
    prstGeom.set("prst", "straightConnector1")
    etree.SubElement(prstGeom, qn("a:avLst"))

    # Line properties
    ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(weight_pt * 12700)))  # pt to EMU

    solidFill = etree.SubElement(ln, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", color_hex.lstrip("#"))

    if style == "dashed":
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")

    return {"type": "connector", "style": style}


def _add_process_arrow(
    slide,
    steps: list[str] = None,
    color_hex: str = "#0F69AF",
    position: str | dict = "below-title",
    slide_width: int = 0,
    slide_height: int = 0,
) -> dict:
    """Add a horizontal process flow with chevron arrows."""
    if not steps:
        steps = ["Step 1", "Step 2", "Step 3"]

    num_steps = len(steps)
    total_width_inches = min(9.0, num_steps * 2.2)
    step_width = Inches(total_width_inches / num_steps)
    step_height = Inches(0.7)
    gap = Inches(0.1)

    # Center the process flow
    total_emu = step_width * num_steps + gap * (num_steps - 1)
    start_left = (slide_width - total_emu) // 2

    _, top = _resolve_position(position, slide_width, slide_height, total_emu, step_height)

    rgb = _hex_to_rgb(color_hex)

    for i, step_text in enumerate(steps):
        left = start_left + i * (step_width + gap)

        # Use chevron for all but first step (which uses a rectangle)
        if i == 0:
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
        else:
            shape_type = MSO_SHAPE.CHEVRON

        shape = slide.shapes.add_shape(shape_type, left, top, step_width, step_height)
        shape.fill.solid()

        # Alternate between full color and lighter shade
        if i % 2 == 0:
            shape.fill.fore_color.rgb = rgb
            fill_hex = color_hex
        else:
            # Lighter shade: blend toward white
            r = min(255, rgb[0] + (255 - rgb[0]) // 2)
            g = min(255, rgb[1] + (255 - rgb[1]) // 2)
            b = min(255, rgb[2] + (255 - rgb[2]) // 2)
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            fill_hex = f"#{r:02X}{g:02X}{b:02X}"

        text_color = _choose_text_color(fill_hex)

        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_text
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

    return {"type": "process_arrow", "steps": steps}


# ---------------------------------------------------------------------------
# Main annotator
# ---------------------------------------------------------------------------


class ShapeAnnotator:
    """Add shape annotations to slides using template theme colors."""

    def __init__(
        self,
        theme_colors: Optional[ThemeColors] = None,
        slide_width: int = 0,
        slide_height: int = 0,
    ):
        self.theme = theme_colors or ThemeColors.uptimize_defaults()
        self.slide_width = slide_width
        self.slide_height = slide_height

    def annotate(self, slide, shapes_defs: list[dict]) -> list[dict]:
        """Add shapes to a slide from a list of shape definitions.

        Args:
            slide: python-pptx Slide object.
            shapes_defs: List of shape definition dicts. Each must have 'type'.

        Returns:
            List of result dicts describing what was added.
        """
        results = []

        for shape_def in shapes_defs:
            shape_type = shape_def.get("type", "").lower()
            color_ref = shape_def.get("color", "accent2")
            color_hex = self.theme.resolve_color(color_ref)
            position = shape_def.get("position", "auto")

            common = {
                "color_hex": color_hex,
                "position": position,
                "slide_width": self.slide_width,
                "slide_height": self.slide_height,
            }

            try:
                if shape_type == "arrow":
                    result = _add_arrow(
                        slide,
                        direction=shape_def.get("direction", "right"),
                        label=shape_def.get("label", ""),
                        **common,
                    )
                elif shape_type == "callout":
                    result = _add_callout(
                        slide,
                        text=shape_def.get("text", ""),
                        **common,
                    )
                elif shape_type == "badge":
                    result = _add_badge(
                        slide,
                        text=shape_def.get("text", ""),
                        **common,
                    )
                elif shape_type == "highlight":
                    result = _add_highlight(
                        slide,
                        opacity=shape_def.get("opacity", 0.3),
                        width_inches=shape_def.get("width", 3.0),
                        height_inches=shape_def.get("height", 1.5),
                        **common,
                    )
                elif shape_type == "connector":
                    result = _add_connector(
                        slide,
                        start=shape_def.get("start", "center-left"),
                        end=shape_def.get("end", "center-right"),
                        style=shape_def.get("style", "solid"),
                        color_hex=color_hex,
                        weight_pt=shape_def.get("weight", 2.0),
                        slide_width=self.slide_width,
                        slide_height=self.slide_height,
                    )
                elif shape_type == "process_arrow":
                    result = _add_process_arrow(
                        slide,
                        steps=shape_def.get("steps", ["Step 1", "Step 2", "Step 3"]),
                        **common,
                    )
                else:
                    log.warning("Unknown shape type: %s", shape_type)
                    continue

                results.append(result)

            except Exception as e:
                log.warning("Failed to add %s shape: %s", shape_type, e)

        return results
