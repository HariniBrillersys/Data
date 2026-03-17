"""
Text Overflow Detection and Auto-Fit.

Estimates whether text will fit inside a placeholder and applies
PowerPoint's "shrink text on overflow" (normAutofit) when needed.

How it works:
  1. After filling a placeholder, estimate how much vertical space the
     text needs given the placeholder width and font size.
  2. If the text would overflow, enable normAutofit on the placeholder's
     bodyPr element. PowerPoint will then shrink the font to fit.
  3. As a second safety net, if the estimated overflow is extreme, also
     set a reduced font size directly so the initial render looks correct
     even in viewers that don't support normAutofit (e.g. LibreOffice).

Font size resolution order (PowerPoint inheritance):
  - Run-level font size (explicit on the run)
  - Paragraph-level defRPr
  - Layout-level placeholder definition
  - Slide master defaults (title=20pt, body=14pt)

Title/subtitle placeholders use a stricter threshold (0.85) to prevent
text wrapping, which looks unprofessional on title lines.
"""

from __future__ import annotations

import logging

from lxml import etree
from pptx.oxml.ns import qn
from pptx.util import Pt

from .template_engine import resolve_font_size_from_xml

log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Average character width as a fraction of font size (em).
_AVG_CHAR_WIDTH_EM = 0.52

# Line height as a fraction of font size.
_LINE_HEIGHT_FACTOR = 1.2

# Minimum font size we'll shrink to (in points).
_MIN_FONT_SIZE_PT = 10

# EMU per point
_EMU_PER_PT = 12700

# Default margins in EMU (PowerPoint defaults)
_DEFAULT_MARGIN_LR = 91440  # ~0.1 inch left/right
_DEFAULT_MARGIN_TB = 45720  # ~0.05 inch top/bottom


# ---------------------------------------------------------------------------
# Margin reading from XML
# ---------------------------------------------------------------------------


def _read_margins(placeholder) -> tuple[int, int, int, int]:
    """
    Read actual margins (insets) from a placeholder's bodyPr element.

    Returns:
        Tuple of (left, right, top, bottom) margins in EMU.
        Falls back to PowerPoint defaults for any missing attribute.
    """
    left = _DEFAULT_MARGIN_LR
    right = _DEFAULT_MARGIN_LR
    top = _DEFAULT_MARGIN_TB
    bottom = _DEFAULT_MARGIN_TB

    try:
        txBody = placeholder._element.find(qn("p:txBody"))
        if txBody is None:
            return left, right, top, bottom

        bodyPr = txBody.find(qn("a:bodyPr"))
        if bodyPr is None:
            return left, right, top, bottom

        lIns = bodyPr.get("lIns")
        if lIns is not None:
            left = int(lIns)
        rIns = bodyPr.get("rIns")
        if rIns is not None:
            right = int(rIns)
        tIns = bodyPr.get("tIns")
        if tIns is not None:
            top = int(tIns)
        bIns = bodyPr.get("bIns")
        if bIns is not None:
            bottom = int(bIns)
    except Exception:
        pass

    return left, right, top, bottom


# ---------------------------------------------------------------------------
# Font size resolution
# ---------------------------------------------------------------------------


def resolve_font_size(placeholder, master_default_pt: float = 20.0) -> float:
    """
    Resolve the effective font size for a placeholder.

    Falls back to defaults based on placeholder type if nothing found in XML.

    Args:
        placeholder: python-pptx placeholder object
        master_default_pt: Default font size for titles (kept for backward compatibility)

    Returns:
        Font size in points
    """
    return resolve_font_size_from_xml(placeholder, default_pt=14.0)


# ---------------------------------------------------------------------------
# Overflow estimation
# ---------------------------------------------------------------------------


def estimate_lines_needed(
    text: str,
    box_width_emu: int,
    font_size_pt: float,
    margin_left_emu: int = _DEFAULT_MARGIN_LR,
    margin_right_emu: int = _DEFAULT_MARGIN_LR,
) -> int:
    """
    Estimate how many lines a text string needs inside a box.

    This is an approximation using average character widths.
    """
    usable_width_emu = box_width_emu - margin_left_emu - margin_right_emu
    usable_width_pt = usable_width_emu / _EMU_PER_PT

    char_width_pt = font_size_pt * _AVG_CHAR_WIDTH_EM
    chars_per_line = max(1, int(usable_width_pt / char_width_pt))

    total_lines = 0
    for paragraph in text.split("\n"):
        paragraph = paragraph.strip()
        if not paragraph:
            total_lines += 1
            continue
        lines_for_para = max(1, -(-len(paragraph) // chars_per_line))
        total_lines += lines_for_para

    return total_lines


def check_overflow(
    text: str,
    box_width_emu: int,
    box_height_emu: int,
    font_size_pt: float,
    margin_left_emu: int = _DEFAULT_MARGIN_LR,
    margin_right_emu: int = _DEFAULT_MARGIN_LR,
    margin_top_emu: int = _DEFAULT_MARGIN_TB,
    margin_bottom_emu: int = _DEFAULT_MARGIN_TB,
) -> dict:
    """
    Check if text will overflow a placeholder box.

    Returns:
        Dict with overflows, lines_needed, lines_available, ratio, suggested_font_pt.
    """
    lines_needed = estimate_lines_needed(text, box_width_emu, font_size_pt, margin_left_emu, margin_right_emu)

    usable_height_emu = box_height_emu - margin_top_emu - margin_bottom_emu
    usable_height_pt = usable_height_emu / _EMU_PER_PT

    line_height_pt = font_size_pt * _LINE_HEIGHT_FACTOR
    lines_available = max(1, int(usable_height_pt / line_height_pt))

    ratio = lines_needed / lines_available
    overflows = ratio > 1.0

    suggested = None
    if overflows:
        lo, hi = _MIN_FONT_SIZE_PT, font_size_pt
        for _ in range(20):
            mid = (lo + hi) / 2
            test_lines = estimate_lines_needed(text, box_width_emu, mid, margin_left_emu, margin_right_emu)
            test_line_height = mid * _LINE_HEIGHT_FACTOR
            test_avail = max(1, int(usable_height_pt / test_line_height))
            if test_lines <= test_avail:
                lo = mid
            else:
                hi = mid
        suggested = max(_MIN_FONT_SIZE_PT, int(lo * 2) / 2)

    return {
        "overflows": overflows,
        "lines_needed": lines_needed,
        "lines_available": lines_available,
        "ratio": round(ratio, 2),
        "current_font_pt": font_size_pt,
        "suggested_font_pt": suggested,
    }


# ---------------------------------------------------------------------------
# Auto-fit application
# ---------------------------------------------------------------------------


def enable_autofit(placeholder) -> None:
    """
    Enable PowerPoint's 'Shrink text on overflow' (normAutofit) on a placeholder.
    """
    txBody = placeholder._element.find(qn("p:txBody"))
    if txBody is None:
        return

    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        return

    for child in list(bodyPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("noAutofit", "normAutofit", "spAutoFit"):
            bodyPr.remove(child)

    etree.SubElement(bodyPr, qn("a:normAutofit"))


def set_font_size_all_runs(placeholder, font_size_pt: float) -> None:
    """Set font size on all runs in a placeholder (fallback for non-PowerPoint viewers)."""
    if not hasattr(placeholder, "text_frame"):
        return

    for para in placeholder.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_size_pt)


# ---------------------------------------------------------------------------
# Post-generation validation
# ---------------------------------------------------------------------------


def _is_title_placeholder(placeholder) -> bool:
    """Check if a placeholder is a title or subtitle type."""
    ph_type = str(placeholder.placeholder_format.type)
    return "TITLE" in ph_type or "SUBTITLE" in ph_type


def validate_slide(slide) -> list[dict]:
    """
    Validate all text placeholders on a slide for overflow.

    Returns a list of issues found (empty list = all good).
    Uses stricter thresholds for title/subtitle placeholders to prevent
    text wrapping on headline elements.
    """
    issues = []

    for ph in slide.placeholders:
        if not hasattr(ph, "text_frame"):
            continue

        text = ph.text.strip()
        if not text:
            continue

        width = ph.width or 0
        height = ph.height or 0
        if width == 0 or height == 0:
            continue

        # Read actual margins from the placeholder XML
        m_left, m_right, m_top, m_bottom = _read_margins(ph)

        font_pt = resolve_font_size(ph)

        result = check_overflow(text, width, height, font_pt, m_left, m_right, m_top, m_bottom)

        # Title/subtitle: use stricter threshold to prevent wrapping
        is_title = _is_title_placeholder(ph)
        overflow_threshold = 0.85 if is_title else 1.0
        is_overflow = result["ratio"] > overflow_threshold

        if is_overflow:
            idx = ph.placeholder_format.idx
            suggested = result["suggested_font_pt"]

            # Enable autofit as a baseline
            enable_autofit(ph)

            if suggested and suggested >= _MIN_FONT_SIZE_PT:
                if result["ratio"] > 1.5:
                    set_font_size_all_runs(ph, suggested)
                    action = f"shrink-to-fit + reduced font to {suggested}pt"
                else:
                    # For moderate overflow (including title 0.85 threshold),
                    # normAutofit alone is sufficient — PowerPoint will shrink
                    # if needed without permanently overwriting the template's
                    # designed font size.
                    action = "enabled shrink-to-fit"
                needs_reword = False
            else:
                set_font_size_all_runs(ph, _MIN_FONT_SIZE_PT)
                suggested = _MIN_FONT_SIZE_PT
                usable_h_emu = (ph.height or 0) - m_top - m_bottom
                usable_h_pt = usable_h_emu / _EMU_PER_PT
                line_h = _MIN_FONT_SIZE_PT * _LINE_HEIGHT_FACTOR
                max_lines = max(1, int(usable_h_pt / line_h))
                usable_w_emu = (ph.width or 0) - m_left - m_right
                usable_w_pt = usable_w_emu / _EMU_PER_PT
                chars_per_line = max(1, int(usable_w_pt / (_MIN_FONT_SIZE_PT * _AVG_CHAR_WIDTH_EM)))
                max_chars = max_lines * chars_per_line
                action = (
                    f"WARNING: text too long even at {_MIN_FONT_SIZE_PT}pt. "
                    f"Shorten to ~{max_chars} chars (currently {len(text)})."
                )
                needs_reword = True

            issues.append(
                {
                    "placeholder_idx": idx,
                    "placeholder_name": ph.name,
                    "text_preview": text[:60] + ("..." if len(text) > 60 else ""),
                    "overflow_ratio": result["ratio"],
                    "lines_needed": result["lines_needed"],
                    "lines_available": result["lines_available"],
                    "original_font_pt": font_pt,
                    "suggested_font_pt": suggested,
                    "action_taken": action,
                    "needs_reword": needs_reword,
                }
            )

    return issues


def validate_presentation(presentation) -> list[dict]:
    """
    Validate all slides in a presentation for text overflow.

    Automatically applies fixes (autofit / font reduction) to any overflows.
    """
    all_issues = []

    for i, slide in enumerate(presentation.slides):
        slide_issues = validate_slide(slide)
        if slide_issues:
            all_issues.append(
                {
                    "slide_number": i + 1,
                    "layout": slide.slide_layout.name,
                    "issues": slide_issues,
                }
            )

    return all_issues
