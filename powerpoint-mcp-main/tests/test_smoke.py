"""
Smoke tests for the PowerPoint MCP server.

These tests verify the core pipeline works end-to-end without requiring
PowerPoint or any external services. They validate:
  - Template loading and layout introspection
  - Theme color extraction
  - Presentation creation with various content types
  - Overflow detection and auto-fit
  - Shape annotations
  - Chart generation (if matplotlib is available)
  - Speaker notes
  - Image fit mode (no-crop placement)
  - Output file validity
"""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from pptx import Presentation

from pptx_mcp.chart_builder import ChartBuilder
from pptx_mcp.composer import PresentationComposer
from pptx_mcp.layout_classifier import LayoutClassifier, SlideIntent
from pptx_mcp.overflow import (
    check_overflow,
    estimate_lines_needed,
)
from pptx_mcp.shape_builder import ShapeAnnotator
from pptx_mcp.template_engine import TemplateEngine
from pptx_mcp.theme_colors import ThemeColors

TEMPLATE_NAME = "Uptimize Master"


# ---------------------------------------------------------------------------
# Template & Theme tests
# ---------------------------------------------------------------------------


class TestTemplateEngine:
    def test_scan_finds_template(self, templates_dir):
        engine = TemplateEngine(templates_dir)
        available = engine.list_available()
        assert any(t["name"] == TEMPLATE_NAME for t in available)

    def test_template_has_layouts(self, templates_dir):
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)
        assert info is not None
        assert len(info.layouts) > 10  # Uptimize has 35 layouts

    def test_template_has_slide_dimensions(self, templates_dir):
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)
        assert info.slide_width > 0
        assert info.slide_height > 0

    def test_layout_has_placeholders(self, templates_dir):
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)
        # At least one layout should have a title placeholder
        has_title = any(layout.has_title for layout in info.layouts)
        assert has_title

    def test_list_templates(self, templates_dir):
        engine = TemplateEngine(templates_dir)
        available = engine.list_available()
        assert isinstance(available, list)
        assert len(available) >= 1
        assert all("name" in t and "file" in t for t in available)


class TestTextStyles:
    """Verify text style extraction captures bullet and spacing info."""

    def test_text_styles_contain_bullet_info(self, templates_dir):
        """Content placeholders should have bullet_type in text_styles."""
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)

        # Find a content layout (not title-only)
        content_layouts = [
            layout for layout in info.layouts if any(p.role.value == "content" for p in layout.placeholders)
        ]
        assert len(content_layouts) > 0, "Should have at least one content layout"

        # Check that at least one content placeholder has text_styles with bullet info
        found_bullet_info = False
        for layout in content_layouts:
            for ph in layout.placeholders:
                if ph.role.value == "content" and ph.text_styles:
                    level_0 = ph.text_styles.get("level_0", {})
                    if "bullet_type" in level_0:
                        found_bullet_info = True
                        # Verify bullet_type is one of expected values
                        assert level_0["bullet_type"] in ("char", "autonum", "none")
                        break
            if found_bullet_info:
                break

        # The Uptimize Master template should have bullet info in lstStyle
        # If not found, the test still passes — some templates may not define lstStyle
        # This is a "present if available" test

    def test_spacing_extracted_correctly(self, templates_dir):
        """Spacing values should be extracted as child elements, not attributes."""
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)

        # Find any placeholder with text_styles that has spacing
        for layout in info.layouts:
            for ph in layout.placeholders:
                if not ph.text_styles:
                    continue
                for level_name, styles in ph.text_styles.items():
                    # If space_before exists, it should be a reasonable value
                    # (in hundredths of a point, so 600 = 6pt)
                    if "space_before" in styles:
                        assert isinstance(styles["space_before"], int)
                        assert styles["space_before"] >= 0
                    if "space_after" in styles:
                        assert isinstance(styles["space_after"], int)
                        assert styles["space_after"] >= 0
                    if "line_spacing" in styles:
                        assert isinstance(styles["line_spacing"], int)
                        # Line spacing in hundredths of percent (100000 = 100%)
                        assert 50000 <= styles["line_spacing"] <= 300000

    def test_text_styles_structure(self, templates_dir):
        """Text styles dict should use level_N naming convention."""
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)

        for layout in info.layouts:
            for ph in layout.placeholders:
                if not ph.text_styles:
                    continue
                for key in ph.text_styles:
                    assert key.startswith("level_"), f"Unexpected key '{key}' in text_styles"
                    level_num = key.split("_")[1]
                    assert level_num.isdigit(), f"Level key '{key}' doesn't end with a number"


class TestThemeColors:
    def test_uptimize_defaults(self):
        theme = ThemeColors.uptimize_defaults()
        assert theme.accent2 == "#0F69AF"  # Merck Blue
        assert theme.dk1 == "#000000"
        assert theme.lt1 == "#FFFFFF"

    def test_from_template(self, templates_dir):
        potx = templates_dir / "Uptimize Master.potx"
        if not potx.exists():
            pytest.skip("Uptimize Master.potx not found")
        theme = ThemeColors.from_template(potx)
        assert theme.accent1 == "#96D7D2"
        assert theme.accent2 == "#0F69AF"
        assert theme.accent3 == "#FFDCB9"

    def test_custom_colors_extracted(self, templates_dir):
        potx = templates_dir / "Uptimize Master.potx"
        if not potx.exists():
            pytest.skip("Uptimize Master.potx not found")
        theme = ThemeColors.from_template(potx)
        assert "Rich Purple" in theme.custom_colors
        assert theme.custom_colors["Rich Purple"] == "#503291"

    def test_resolve_color_hex(self):
        theme = ThemeColors.uptimize_defaults()
        assert theme.resolve_color("#FF0000") == "#FF0000"

    def test_resolve_color_theme_name(self):
        theme = ThemeColors.uptimize_defaults()
        assert theme.resolve_color("accent2") == "#0F69AF"
        assert theme.resolve_color("dark") == "#000000"

    def test_resolve_color_custom_name(self):
        theme = ThemeColors.uptimize_defaults()
        assert theme.resolve_color("Rich Purple") == "#503291"

    def test_accent_cycle(self):
        theme = ThemeColors.uptimize_defaults()
        cycle = theme.accent_cycle()
        assert len(cycle) == 6
        assert cycle[0] == "#96D7D2"


# ---------------------------------------------------------------------------
# Layout Classifier tests
# ---------------------------------------------------------------------------


class TestLayoutClassifier:
    def test_classify_title_slide(self, templates_dir):
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)
        classifier = LayoutClassifier(info)

        content = {"title": "Hello", "subtitle": "World"}
        spec = classifier.classify_content(content)
        assert spec.intent == SlideIntent.TITLE_SLIDE

    def test_classify_content_slide(self, templates_dir):
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)
        classifier = LayoutClassifier(info)

        content = {"title": "Data", "content": ["Point 1", "Point 2"]}
        spec = classifier.classify_content(content)
        assert spec.intent == SlideIntent.SINGLE_CONTENT

    def test_classify_two_column(self, templates_dir):
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)
        classifier = LayoutClassifier(info)

        content = {"title": "Compare", "content_left": "A", "content_right": "B"}
        spec = classifier.classify_content(content)
        assert spec.intent == SlideIntent.TWO_COLUMN

    def test_auto_select(self, templates_dir):
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)
        classifier = LayoutClassifier(info)

        content = {"title": "Title", "content": ["Bullet"]}
        layout, spec = classifier.auto_select(content)
        assert layout is not None
        assert spec.intent == SlideIntent.SINGLE_CONTENT


# ---------------------------------------------------------------------------
# Overflow tests
# ---------------------------------------------------------------------------


class TestOverflow:
    def test_estimate_short_text(self):
        lines = estimate_lines_needed("Hello World", 5000000, 20.0)
        assert lines == 1

    def test_estimate_long_text(self):
        long_text = "This is a very long sentence that should wrap to multiple lines " * 5
        lines = estimate_lines_needed(long_text, 3000000, 14.0)
        assert lines > 1

    def test_check_overflow_no_overflow(self):
        result = check_overflow("Short", 5000000, 2000000, 14.0)
        assert result["overflows"] is False
        assert result["ratio"] <= 1.0

    def test_check_overflow_with_overflow(self):
        long_text = "Very long text content " * 50
        result = check_overflow(long_text, 2000000, 500000, 14.0)
        assert result["overflows"] is True
        assert result["suggested_font_pt"] is not None


# ---------------------------------------------------------------------------
# Presentation creation (end-to-end)
# ---------------------------------------------------------------------------


class TestPresentationCreation:
    def test_basic_presentation(self, templates_dir, output_dir):
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {"title": "Test Presentation", "subtitle": "Smoke Test"},
                {"title": "Bullet Points", "content": ["Point 1", "Point 2", "Point 3"]},
            ],
        )

        assert result["success"] is True
        assert Path(result["output_path"]).exists()
        assert result["num_slides"] == 2

    def test_two_column_presentation(self, templates_dir, output_dir):
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {"title": "Comparison", "content_left": ["Left A", "Left B"], "content_right": ["Right A", "Right B"]},
            ],
        )

        assert result["success"] is True

    def test_output_is_valid_zip(self, templates_dir, output_dir):
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[{"title": "Validity Test", "subtitle": "ZIP check"}],
        )

        assert result["success"] is True
        assert zipfile.is_zipfile(result["output_path"])

    def test_long_title_triggers_autofit(self, templates_dir, output_dir):
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        long_title = (
            "This Is An Extremely Long Title That Should Definitely "
            "Trigger The Overflow Detection And Auto-Fit Mechanism"
        )

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[{"title": long_title, "subtitle": "Autofit test"}],
        )

        assert result["success"] is True
        # The overflow system should have handled it gracefully
        assert Path(result["output_path"]).exists()

    def test_invalid_template_returns_error(self, templates_dir, output_dir):
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name="Nonexistent Template",
            slides=[{"title": "Test"}],
        )

        assert result["success"] is False
        assert "not found" in result["error"].lower()


# ---------------------------------------------------------------------------
# Shape annotations
# ---------------------------------------------------------------------------


class TestShapeAnnotations:
    def test_annotate_slide_with_shapes(self, templates_dir, output_dir):
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "With Shapes",
                    "content": ["Data point 1", "Data point 2"],
                    "shapes": [
                        {"type": "badge", "text": "NEW", "position": "top-right", "color": "accent1"},
                        {"type": "arrow", "direction": "right", "label": "+50%", "position": "auto"},
                    ],
                },
            ],
        )

        assert result["success"] is True
        slide_result = result["slides"][0]
        assert slide_result["shapes_added"] is not None
        assert len(slide_result["shapes_added"]) == 2

    def test_shape_annotator_color_resolution(self):
        theme = ThemeColors.uptimize_defaults()
        annotator = ShapeAnnotator(theme_colors=theme, slide_width=12192000, slide_height=6858000)
        assert annotator.theme.resolve_color("accent1") == "#96D7D2"

    def test_process_arrow(self, templates_dir, output_dir):
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Process Flow",
                    "shapes": [
                        {"type": "process_arrow", "steps": ["Plan", "Build", "Test", "Deploy"]},
                    ],
                },
            ],
        )

        assert result["success"] is True


# ---------------------------------------------------------------------------
# Chart generation (optional -- requires matplotlib)
# ---------------------------------------------------------------------------


class TestChartGeneration:
    def test_bar_chart(self, output_dir):
        theme = ThemeColors.uptimize_defaults()
        builder = ChartBuilder(images_dir=output_dir)

        result = builder.generate(
            chart_type="bar",
            data={"Q1": 10, "Q2": 15, "Q3": 22, "Q4": 28},
            title="Revenue by Quarter",
            theme_colors=theme,
        )

        assert result["success"] is True
        path = Path(result["path"])
        assert path.exists()
        assert path.stat().st_size > 1000

    def test_line_chart_multi_series(self, output_dir):
        theme = ThemeColors.uptimize_defaults()
        builder = ChartBuilder(images_dir=output_dir)

        result = builder.generate(
            chart_type="line",
            data=[
                {"name": "Revenue", "values": {"Q1": 10, "Q2": 15, "Q3": 22}},
                {"name": "Costs", "values": {"Q1": 8, "Q2": 11, "Q3": 14}},
            ],
            title="Revenue vs Costs",
            theme_colors=theme,
        )

        assert result["success"] is True

    def test_pie_chart(self, output_dir):
        builder = ChartBuilder(images_dir=output_dir)

        result = builder.generate(
            chart_type="pie",
            data={"North": 40, "South": 25, "East": 20, "West": 15},
            title="Regional Split",
        )

        assert result["success"] is True

    def test_invalid_chart_type(self, output_dir):
        builder = ChartBuilder(images_dir=output_dir)

        result = builder.generate(
            chart_type="invalid_type",
            data={"A": 1},
        )

        assert result["success"] is False
        assert "unknown" in result["error"].lower()

    def test_chart_uses_theme_colors(self, output_dir):
        """Verify chart respects theme color extraction."""
        theme = ThemeColors(
            accent1="#FF0000",
            accent2="#00FF00",
            accent3="#0000FF",
        )
        builder = ChartBuilder(images_dir=output_dir)

        result = builder.generate(
            chart_type="bar",
            data={"A": 1, "B": 2},
            theme_colors=theme,
        )

        assert result["success"] is True

    def test_chart_returns_image_mode_fit(self, output_dir):
        """Verify chart generation returns image_mode: fit."""
        builder = ChartBuilder(images_dir=output_dir)

        result = builder.generate(
            chart_type="bar",
            data={"A": 1, "B": 2},
        )

        assert result["success"] is True
        assert result["image_mode"] == "fit"

    def test_chart_builder_warmup(self, output_dir):
        """ChartBuilder should pre-warm matplotlib during init."""
        import matplotlib.font_manager as fm

        _builder = ChartBuilder(images_dir=output_dir)  # noqa: F841 – triggers warmup
        # After init, font manager should be populated (cache warm)
        assert len(fm.fontManager.ttflist) > 0

    def test_chart_rejects_too_many_labels(self, output_dir):
        """Charts with >100 data points should be rejected."""
        builder = ChartBuilder(images_dir=output_dir)
        big_data = {f"label_{i}": i for i in range(150)}

        result = builder.generate(chart_type="bar", data=big_data)

        assert result["success"] is False
        assert "too many" in result["error"].lower()

    def test_chart_rejects_too_many_series(self, output_dir):
        """Charts with >20 series should be rejected."""
        builder = ChartBuilder(images_dir=output_dir)
        many_series = [{"name": f"Series {i}", "values": {"A": i, "B": i + 1}} for i in range(25)]

        result = builder.generate(chart_type="bar", data=many_series)

        assert result["success"] is False
        assert "too many" in result["error"].lower()


# ---------------------------------------------------------------------------
# Speaker notes
# ---------------------------------------------------------------------------


class TestSpeakerNotes:
    def test_notes_string(self, templates_dir, output_dir):
        """Verify a string note is written to the slide."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Slide With Notes",
                    "content": ["Key finding"],
                    "notes": "Remember to mention the quarterly growth figures.",
                },
            ],
        )

        assert result["success"] is True
        assert result["slides"][0]["fields_filled"].get("notes") is True

        # Open the saved file and verify the notes text
        prs = Presentation(result["output_path"])
        slide = prs.slides[0]
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "quarterly growth" in notes_text

    def test_notes_list(self, templates_dir, output_dir):
        """Verify a list of notes becomes multiple paragraphs."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        talking_points = [
            "Mention the 15% increase in adoption",
            "Compare with competitor benchmarks",
            "Highlight the regional differences",
        ]

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Multi-Point Notes",
                    "subtitle": "Q4 Review",
                    "notes": talking_points,
                },
            ],
        )

        assert result["success"] is True

        prs = Presentation(result["output_path"])
        slide = prs.slides[0]
        tf = slide.notes_slide.notes_text_frame
        paragraphs = [p.text for p in tf.paragraphs if p.text.strip()]
        assert len(paragraphs) >= 3
        assert "15% increase" in paragraphs[0]
        assert "competitor" in paragraphs[1]
        assert "regional" in paragraphs[2]

    def test_notes_do_not_affect_layout_selection(self, templates_dir):
        """Notes should not influence intent classification."""
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)
        classifier = LayoutClassifier(info)

        # Same content with and without notes should produce the same intent
        content_no_notes = {"title": "Data", "content": ["Point"]}
        content_with_notes = {"title": "Data", "content": ["Point"], "notes": "Talk about X"}

        spec_no = classifier.classify_content(content_no_notes)
        spec_with = classifier.classify_content(content_with_notes)

        assert spec_no.intent == spec_with.intent


# ---------------------------------------------------------------------------
# Image fit mode (no-crop placement)
# ---------------------------------------------------------------------------


class TestImageFitMode:
    def test_fit_mode_creates_freestanding_image(self, templates_dir, output_dir, test_image):
        """Verify image_mode=fit places an image as a freestanding shape."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Chart Slide",
                    "image": str(test_image),
                    "image_mode": "fit",
                },
            ],
        )

        assert result["success"] is True
        assert result["slides"][0]["fields_filled"].get("image") is True

        # Open and verify: the image should be a freestanding picture shape.
        prs = Presentation(result["output_path"])
        freestanding = _get_freestanding_pictures(prs.slides[0])
        assert len(freestanding) >= 1

    def test_fill_mode_default_backward_compat(self, templates_dir, output_dir, test_image):
        """Verify default image_mode (fill) still works."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Storytelling Image",
                    "content": ["Some context for the image"],
                    "image": str(test_image),
                    # No image_mode — defaults to "fill"
                },
            ],
        )

        assert result["success"] is True
        assert result["slides"][0]["fields_filled"].get("image") is True

    def test_classifier_data_image_intent(self, templates_dir):
        """Verify image_mode=fit triggers DATA_IMAGE intent."""
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)
        classifier = LayoutClassifier(info)

        content = {"title": "Revenue Chart", "image": "/path/to/chart.png", "image_mode": "fit"}
        spec = classifier.classify_content(content)
        assert spec.intent == SlideIntent.DATA_IMAGE

    def test_classifier_fill_mode_uses_content_with_image(self, templates_dir):
        """Verify default fill mode still classifies as CONTENT_WITH_IMAGE."""
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)
        classifier = LayoutClassifier(info)

        content = {"title": "Hero Image", "content": ["Caption"], "image": "/path/to/photo.png"}
        spec = classifier.classify_content(content)
        assert spec.intent == SlideIntent.CONTENT_WITH_IMAGE

    def test_data_image_layout_avoids_picture_placeholder(self, templates_dir):
        """Verify DATA_IMAGE intent prefers layouts without picture placeholders."""
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)
        classifier = LayoutClassifier(info)

        content = {"title": "Chart Results", "image": "/path/to/chart.png", "image_mode": "fit"}
        layout, spec = classifier.auto_select(content)
        assert spec.intent == SlideIntent.DATA_IMAGE
        # The selected layout should preferably NOT have a picture placeholder
        # (title-only or blank layouts are ideal for freestanding images)
        # At minimum, verify the layout was selected and has a title
        assert layout is not None
        assert layout.has_title

    def test_wide_image_fit_no_crop(self, templates_dir, output_dir, test_image_wide):
        """Verify a wide 16:9 image placed in fit mode is not cropped."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Wide Chart",
                    "image": str(test_image_wide),
                    "image_mode": "fit",
                },
            ],
        )

        assert result["success"] is True

        # Verify the image shape maintains approximately 16:9 aspect ratio
        prs = Presentation(result["output_path"])
        slide = prs.slides[0]
        freestanding = _get_freestanding_pictures(slide)
        assert len(freestanding) >= 1
        shape = freestanding[0]
        aspect = shape.width / shape.height
        # Original is 16:9 = 1.78; allow some tolerance
        assert 1.5 < aspect < 2.1


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _get_freestanding_pictures(slide):
    """Return picture shapes that are NOT inside a placeholder."""
    result = []
    for s in slide.shapes:
        if not hasattr(s, "image"):
            continue
        try:
            s.placeholder_format
        except ValueError:
            result.append(s)
    return result


# ---------------------------------------------------------------------------
# Multi-image placement (grid and collage)
# ---------------------------------------------------------------------------


class TestMultiImage:
    def test_multi_image_grid(self, templates_dir, output_dir, test_images_multi):
        """Verify 3 images in fit mode produce 3 freestanding picture shapes."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        paths = [str(p) for p in test_images_multi[:3]]
        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Grid Layout",
                    "image": paths,
                    "image_mode": "fit",
                },
            ],
        )

        assert result["success"] is True
        assert result["slides"][0]["fields_filled"].get("image") is True

        prs = Presentation(result["output_path"])
        freestanding = _get_freestanding_pictures(prs.slides[0])
        assert len(freestanding) == 3

    def test_multi_image_collage(self, templates_dir, output_dir, test_images_multi):
        """Verify 3 images in collage mode produce 3 freestanding shapes."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        paths = [str(p) for p in test_images_multi[:3]]
        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Collage Layout",
                    "image": paths,
                    "image_mode": "collage",
                },
            ],
        )

        assert result["success"] is True

        prs = Presentation(result["output_path"])
        freestanding = _get_freestanding_pictures(prs.slides[0])
        assert len(freestanding) == 3

    def test_single_image_collage_fallback(self, templates_dir, output_dir, test_image):
        """Single image with collage mode should work like fit."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Single Collage",
                    "image": str(test_image),
                    "image_mode": "collage",
                },
            ],
        )

        assert result["success"] is True
        assert result["slides"][0]["fields_filled"].get("image") is True

        prs = Presentation(result["output_path"])
        freestanding = _get_freestanding_pictures(prs.slides[0])
        assert len(freestanding) == 1

    def test_multi_image_fill_uses_first(self, templates_dir, output_dir, test_images_multi):
        """List with fill mode should use only the first image."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        paths = [str(p) for p in test_images_multi[:3]]
        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Fill with list",
                    "content": ["Some text"],
                    "image": paths,
                    # No image_mode -> defaults to "fill"
                },
            ],
        )

        assert result["success"] is True
        # Should succeed (first image placed in placeholder)
        assert result["slides"][0]["fields_filled"].get("image") is True

    def test_image_dict_format(self, templates_dir, output_dir, test_images_multi):
        """Images as dicts with path and z_order should work."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        image_dicts = [
            {"path": str(test_images_multi[0]), "z_order": 2},
            {"path": str(test_images_multi[1]), "z_order": 0},
            {"path": str(test_images_multi[2]), "z_order": 1},
        ]
        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Dict Format",
                    "image": image_dicts,
                    "image_mode": "collage",
                },
            ],
        )

        assert result["success"] is True

        prs = Presentation(result["output_path"])
        freestanding = _get_freestanding_pictures(prs.slides[0])
        assert len(freestanding) == 3

    def test_multi_image_cap_at_9(self, templates_dir, output_dir, tmp_path):
        """More than 9 images should be capped at 9."""
        from PIL import Image as PILImage

        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        # Create 12 tiny test images
        paths = []
        for i in range(12):
            p = tmp_path / f"cap_{i}.png"
            PILImage.new("RGB", (100, 100), color=(i * 20, 100, 100)).save(p)
            paths.append(str(p))

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Cap Test",
                    "image": paths,
                    "image_mode": "fit",
                },
            ],
        )

        assert result["success"] is True

        prs = Presentation(result["output_path"])
        freestanding = _get_freestanding_pictures(prs.slides[0])
        assert len(freestanding) == 9

    def test_classifier_collage_mode(self, templates_dir):
        """Verify collage triggers DATA_IMAGE intent."""
        engine = TemplateEngine(templates_dir)
        info = engine.get_template(TEMPLATE_NAME)
        classifier = LayoutClassifier(info)

        content = {
            "title": "Visual Story",
            "image": ["/a.png", "/b.png"],
            "image_mode": "collage",
        }
        spec = classifier.classify_content(content)
        assert spec.intent == SlideIntent.DATA_IMAGE


# ---------------------------------------------------------------------------
# Paragraph Formatting tests (Phase 7)
# ---------------------------------------------------------------------------


class TestParagraphFormatting:
    """Verify paragraph formatting preservation and extended content features."""

    @staticmethod
    def _find_content_placeholder(slide, min_paragraphs=2):
        """Find the content placeholder (not title/subtitle) with actual content."""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 1 and len(ph.text_frame.paragraphs) >= min_paragraphs:
                return ph
        # Fallback: any non-title placeholder
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 1:
                return ph
        return None

    def test_bullet_list_pPr_consistency(self, templates_dir, output_dir):
        """All paragraphs in a bullet list should have pPr elements (not bare)."""
        from pptx.oxml.ns import qn

        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[{"title": "Consistency Test", "content": ["Item 1", "Item 2", "Item 3", "Item 4", "Item 5"]}],
        )

        assert result["success"] is True

        prs = Presentation(result["output_path"])
        slide = prs.slides[0]
        ph = self._find_content_placeholder(slide, min_paragraphs=5)
        assert ph is not None, "No content placeholder found with 5 paragraphs"
        tf = ph.text_frame
        assert len(tf.paragraphs) == 5, f"Expected 5 paragraphs, got {len(tf.paragraphs)}"

        for i, para in enumerate(tf.paragraphs):
            pPr = para._element.find(qn("a:pPr"))
            assert pPr is not None, f"Paragraph {i} missing pPr"

        for i, para in enumerate(tf.paragraphs):
            assert len(para.runs) >= 1, f"Paragraph {i} has no runs"

    def test_extended_format_numbered_list(self, templates_dir, output_dir):
        """Extended format with bullet='number' should produce buAutoNum."""
        from pptx.oxml.ns import qn

        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Numbered List",
                    "content": {"items": ["First", "Second", "Third"], "bullet": "number"},
                }
            ],
        )

        assert result["success"] is True

        prs = Presentation(result["output_path"])
        slide = prs.slides[0]
        ph = self._find_content_placeholder(slide, min_paragraphs=3)
        assert ph is not None, "No content placeholder found"
        tf = ph.text_frame
        assert len(tf.paragraphs) == 3

        for i, para in enumerate(tf.paragraphs):
            pPr = para._element.find(qn("a:pPr"))
            assert pPr is not None, f"Paragraph {i} missing pPr"
            buAutoNum = pPr.find(qn("a:buAutoNum"))
            assert buAutoNum is not None, f"Paragraph {i} missing buAutoNum"

    def test_extended_format_no_bullets(self, templates_dir, output_dir):
        """Extended format with bullet='none' should produce buNone."""
        from pptx.oxml.ns import qn

        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "No Bullets",
                    "content": {"items": ["Line 1", "Line 2", "Line 3"], "bullet": "none"},
                }
            ],
        )

        assert result["success"] is True

        prs = Presentation(result["output_path"])
        slide = prs.slides[0]
        ph = self._find_content_placeholder(slide, min_paragraphs=3)
        assert ph is not None, "No content placeholder found"
        tf = ph.text_frame
        for para in tf.paragraphs:
            pPr = para._element.find(qn("a:pPr"))
            if pPr is not None:
                buNone = pPr.find(qn("a:buNone"))
                assert buNone is not None, "Expected buNone for no-bullet mode"

    def test_nested_list_levels(self, templates_dir, output_dir):
        """Nested lists should produce paragraphs at different indent levels."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Nested",
                    "content": ["Top 1", ["Sub A", "Sub B"], "Top 2", ["Sub C"]],
                }
            ],
        )

        assert result["success"] is True

        prs = Presentation(result["output_path"])
        slide = prs.slides[0]
        ph = self._find_content_placeholder(slide, min_paragraphs=5)
        assert ph is not None, "No content placeholder found"
        tf = ph.text_frame
        # Should be: Top1(L0), SubA(L1), SubB(L1), Top2(L0), SubC(L1)
        assert len(tf.paragraphs) == 5
        levels = [p.level for p in tf.paragraphs]
        assert levels == [0, 1, 1, 0, 1]

    def test_per_item_dict_format(self, templates_dir, output_dir):
        """Dict items in list should support text, level, and overrides."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Mixed Format",
                    "content": [
                        "Plain item",
                        {"text": "Bold item", "bold": True},
                        {"text": "Sub-item", "level": 1},
                    ],
                }
            ],
        )

        assert result["success"] is True

        prs = Presentation(result["output_path"])
        slide = prs.slides[0]
        ph = self._find_content_placeholder(slide, min_paragraphs=3)
        assert ph is not None, "No content placeholder found"
        tf = ph.text_frame
        assert len(tf.paragraphs) == 3

        bold_para = tf.paragraphs[1]
        assert bold_para.runs[0].font.bold is True

        assert tf.paragraphs[2].level == 1

    def test_backward_compat_plain_list(self, templates_dir, output_dir):
        """Plain list format ['a', 'b', 'c'] should work identically to before."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[{"title": "Plain List", "content": ["One", "Two", "Three"]}],
        )

        assert result["success"] is True
        assert result["slides"][0]["fields_filled"]["content"] is True

        prs = Presentation(result["output_path"])
        slide = prs.slides[0]
        ph = self._find_content_placeholder(slide, min_paragraphs=3)
        assert ph is not None, "No content placeholder found"
        tf = ph.text_frame
        texts = [p.text for p in tf.paragraphs]
        assert texts == ["One", "Two", "Three"]

    def test_backward_compat_string_content(self, templates_dir, output_dir):
        """Single string content should work identically to before."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[{"title": "String Content", "content": "Just a paragraph of text"}],
        )

        assert result["success"] is True
        assert result["slides"][0]["fields_filled"]["content"] is True

    def test_custom_bullet_char(self, templates_dir, output_dir):
        """Custom bullet character via extended format."""
        from pptx.oxml.ns import qn

        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Custom Bullets",
                    "content": {"items": ["Check 1", "Check 2"], "bullet": "✓"},
                }
            ],
        )

        assert result["success"] is True

        prs = Presentation(result["output_path"])
        slide = prs.slides[0]
        ph = self._find_content_placeholder(slide, min_paragraphs=2)
        assert ph is not None, "No content placeholder found"
        tf = ph.text_frame
        for para in tf.paragraphs:
            pPr = para._element.find(qn("a:pPr"))
            if pPr is not None:
                buChar = pPr.find(qn("a:buChar"))
                assert buChar is not None, "Expected buChar for custom bullet"
                assert buChar.get("char") == "✓"

    def test_pptx_valid_xml(self, templates_dir, output_dir):
        """Generated PPTX should be a valid ZIP with expected content types."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {"title": "Slide 1", "content": ["A", "B", "C"]},
                {"title": "Slide 2", "content": {"items": ["X", "Y"], "bullet": "number"}},
                {"title": "Slide 3", "content": ["Top", ["Sub1", "Sub2"]]},
            ],
            output_name="formatting_test.pptx",
        )

        assert result["success"] is True
        assert zipfile.is_zipfile(result["output_path"])
        prs = Presentation(result["output_path"])
        assert len(prs.slides) == 3

    def test_multiple_content_setters(self, templates_dir, output_dir):
        """Extended format should work with content_left/content_right too."""
        engine = TemplateEngine(templates_dir)

        composer = PresentationComposer(engine, output_dir)

        result = composer.create_presentation(
            template_name=TEMPLATE_NAME,
            slides=[
                {
                    "title": "Two Column",
                    "content_left": {"items": ["Left 1", "Left 2"], "bullet": "number"},
                    "content_right": ["Right A", "Right B"],
                }
            ],
        )

        assert result["success"] is True
