"""Comprehensive unit tests for text_formatter module."""

from __future__ import annotations

from pathlib import Path

import pytest
from lxml import etree
from pptx.oxml.ns import qn

from pptx_mcp.text_formatter import (
    BulletSpec,
    apply_bullet_type,
    copy_paragraph_properties,
    copy_run_properties,
    ensure_pPr,
    insert_pPr_child_ordered,
    resolve_level_style,
)

# ============================================================================
# Test Helpers
# ============================================================================


def _make_paragraph(text="Test", level=0, with_pPr=True, bullet_char=None, autonum=None):
    """Create a minimal a:p XML element for testing."""
    p = etree.Element(qn("a:p"))
    if with_pPr:
        pPr = etree.SubElement(p, qn("a:pPr"))
        if level > 0:
            pPr.set("lvl", str(level))
        if bullet_char:
            buChar = etree.SubElement(pPr, qn("a:buChar"))
            buChar.set("char", bullet_char)
        if autonum:
            buAutoNum = etree.SubElement(pPr, qn("a:buAutoNum"))
            buAutoNum.set("type", autonum)
    r = etree.SubElement(p, qn("a:r"))
    rPr = etree.SubElement(r, qn("a:rPr"))
    rPr.set("lang", "en-US")
    t = etree.SubElement(r, qn("a:t"))
    t.text = text
    return p


def _get_child_tags(elem):
    """Get list of local tag names of children."""
    return [child.tag.split("}")[-1] if "}" in child.tag else child.tag for child in elem]


# ============================================================================
# Test Classes
# ============================================================================


class TestInsertPPrChildOrdered:
    """Tests for XML element ordering helper."""

    def test_insert_bullet_before_defRPr(self):
        """pPr has defRPr → insert buChar → buChar appears before defRPr."""
        pPr = etree.Element(qn("a:pPr"))
        etree.SubElement(pPr, qn("a:defRPr"))

        buChar = etree.Element(qn("a:buChar"))
        buChar.set("char", "•")
        insert_pPr_child_ordered(pPr, buChar)

        tags = _get_child_tags(pPr)
        assert tags == ["buChar", "defRPr"]

    def test_insert_spacing_before_bullet(self):
        """pPr has buChar → insert lnSpc → lnSpc appears before buChar."""
        pPr = etree.Element(qn("a:pPr"))
        etree.SubElement(pPr, qn("a:buChar"))

        lnSpc = etree.Element(qn("a:lnSpc"))
        insert_pPr_child_ordered(pPr, lnSpc)

        tags = _get_child_tags(pPr)
        assert tags == ["lnSpc", "buChar"]

    def test_insert_into_empty_pPr(self):
        """empty pPr → insert buAutoNum → element added."""
        pPr = etree.Element(qn("a:pPr"))

        buAutoNum = etree.Element(qn("a:buAutoNum"))
        insert_pPr_child_ordered(pPr, buAutoNum)

        tags = _get_child_tags(pPr)
        assert tags == ["buAutoNum"]

    def test_insert_preserves_existing_order(self):
        """pPr has lnSpc + defRPr → insert buChar → order is lnSpc, buChar, defRPr."""
        pPr = etree.Element(qn("a:pPr"))
        etree.SubElement(pPr, qn("a:lnSpc"))
        etree.SubElement(pPr, qn("a:defRPr"))

        buChar = etree.Element(qn("a:buChar"))
        insert_pPr_child_ordered(pPr, buChar)

        tags = _get_child_tags(pPr)
        assert tags == ["lnSpc", "buChar", "defRPr"]

    def test_insert_defRPr_at_end(self):
        """insert defRPr → goes after all bullet elements."""
        pPr = etree.Element(qn("a:pPr"))
        etree.SubElement(pPr, qn("a:lnSpc"))
        etree.SubElement(pPr, qn("a:buChar"))

        defRPr = etree.Element(qn("a:defRPr"))
        insert_pPr_child_ordered(pPr, defRPr)

        tags = _get_child_tags(pPr)
        assert tags == ["lnSpc", "buChar", "defRPr"]


class TestEnsurePPr:
    """Tests for ensure_pPr helper."""

    def test_creates_pPr_when_missing(self):
        """paragraph without pPr → ensure_pPr → pPr exists as first child."""
        p = _make_paragraph(with_pPr=False)
        initial_children = len(p)

        pPr = ensure_pPr(p)

        assert pPr is not None
        assert pPr.tag == qn("a:pPr")
        assert len(p) == initial_children + 1
        # pPr should be first child
        assert p[0] == pPr

    def test_returns_existing_pPr(self):
        """paragraph with pPr → ensure_pPr → returns same element, no duplicate."""
        p = _make_paragraph(with_pPr=True)
        original_pPr = p.find(qn("a:pPr"))
        initial_children = len(p)

        pPr = ensure_pPr(p)

        assert pPr is original_pPr
        assert len(p) == initial_children  # No new element added

    def test_pPr_before_runs(self):
        """paragraph with runs but no pPr → ensure_pPr → pPr is first child, before a:r."""
        p = _make_paragraph(with_pPr=False)
        # Verify we have run elements
        runs = p.findall(qn("a:r"))
        assert len(runs) > 0

        pPr = ensure_pPr(p)

        # pPr should be first child (index 0)
        assert p[0] == pPr
        # First run should now be at index 1
        assert p[1].tag == qn("a:r")


class TestCopyParagraphProperties:
    """Tests for copy_paragraph_properties."""

    def test_copies_bullet_char(self):
        """source has buChar → copy to bare target → target has buChar with same char."""
        source_p = _make_paragraph(bullet_char="▸")
        source_pPr = source_p.find(qn("a:pPr"))

        target_p = _make_paragraph(with_pPr=False)

        copy_paragraph_properties(source_pPr, target_p)

        target_pPr = target_p.find(qn("a:pPr"))
        assert target_pPr is not None
        buChar = target_pPr.find(qn("a:buChar"))
        assert buChar is not None
        assert buChar.get("char") == "▸"

    def test_copies_spacing(self):
        """source has lnSpc, spcBef, spcAft → copy → target has identical spacing elements."""
        source_p = _make_paragraph()
        source_pPr = source_p.find(qn("a:pPr"))

        # Add spacing elements to source
        lnSpc = etree.SubElement(source_pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", "115000")

        spcBef = etree.SubElement(source_pPr, qn("a:spcBef"))
        spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
        spcPts.set("val", "600")

        spcAft = etree.SubElement(source_pPr, qn("a:spcAft"))
        spcPts2 = etree.SubElement(spcAft, qn("a:spcPts"))
        spcPts2.set("val", "300")

        target_p = _make_paragraph(with_pPr=False)

        copy_paragraph_properties(source_pPr, target_p)

        target_pPr = target_p.find(qn("a:pPr"))
        assert target_pPr is not None

        # Check lnSpc
        lnSpc_copy = target_pPr.find(qn("a:lnSpc"))
        assert lnSpc_copy is not None
        spcPct_copy = lnSpc_copy.find(qn("a:spcPct"))
        assert spcPct_copy is not None
        assert spcPct_copy.get("val") == "115000"

        # Check spcBef
        spcBef_copy = target_pPr.find(qn("a:spcBef"))
        assert spcBef_copy is not None
        spcPts_copy = spcBef_copy.find(qn("a:spcPts"))
        assert spcPts_copy is not None
        assert spcPts_copy.get("val") == "600"

        # Check spcAft
        spcAft_copy = target_pPr.find(qn("a:spcAft"))
        assert spcAft_copy is not None
        spcPts_copy2 = spcAft_copy.find(qn("a:spcPts"))
        assert spcPts_copy2 is not None
        assert spcPts_copy2.get("val") == "300"

    def test_copies_line_spacing_element(self):
        """source pPr has lnSpc>spcPct element → copy → target pPr contains lnSpc child."""
        source_p = _make_paragraph()
        source_pPr = source_p.find(qn("a:pPr"))

        lnSpc = etree.SubElement(source_pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", "115000")

        target_p = _make_paragraph(with_pPr=False)

        copy_paragraph_properties(source_pPr, target_p)

        target_pPr = target_p.find(qn("a:pPr"))
        lnSpc_copy = target_pPr.find(qn("a:lnSpc"))
        assert lnSpc_copy is not None
        spcPct_copy = lnSpc_copy.find(qn("a:spcPct"))
        assert spcPct_copy is not None
        assert spcPct_copy.get("val") == "115000"

    def test_copies_indentation(self):
        """source has marL and indent attributes → copy → target has them."""
        source_p = _make_paragraph()
        source_pPr = source_p.find(qn("a:pPr"))
        source_pPr.set("marL", "457200")
        source_pPr.set("indent", "-228600")

        target_p = _make_paragraph(with_pPr=False)

        copy_paragraph_properties(source_pPr, target_p)

        target_pPr = target_p.find(qn("a:pPr"))
        assert target_pPr.get("marL") == "457200"
        assert target_pPr.get("indent") == "-228600"

    def test_does_not_copy_level(self):
        """source pPr has lvl='1' → copy → target pPr has NO lvl attribute."""
        source_p = _make_paragraph(level=1)
        source_pPr = source_p.find(qn("a:pPr"))
        assert source_pPr.get("lvl") == "1"

        target_p = _make_paragraph(with_pPr=False)

        copy_paragraph_properties(source_pPr, target_p)

        target_pPr = target_p.find(qn("a:pPr"))
        assert target_pPr.get("lvl") is None

    def test_replaces_existing_pPr(self):
        """target already has pPr with different content → copy → target has source's content."""
        source_p = _make_paragraph(bullet_char="▸")
        source_pPr = source_p.find(qn("a:pPr"))

        target_p = _make_paragraph(bullet_char="•")

        copy_paragraph_properties(source_pPr, target_p)

        target_pPr = target_p.find(qn("a:pPr"))
        buChar = target_pPr.find(qn("a:buChar"))
        assert buChar.get("char") == "▸"

    def test_deep_copy_isolation(self):
        """modify copied pPr on target → source pPr unchanged (no reference sharing)."""
        source_p = _make_paragraph(bullet_char="▸")
        source_pPr = source_p.find(qn("a:pPr"))

        target_p = _make_paragraph(with_pPr=False)

        copy_paragraph_properties(source_pPr, target_p)

        # Modify target's pPr
        target_pPr = target_p.find(qn("a:pPr"))
        target_buChar = target_pPr.find(qn("a:buChar"))
        target_buChar.set("char", "NEW")

        # Source should be unchanged
        source_buChar = source_pPr.find(qn("a:buChar"))
        assert source_buChar.get("char") == "▸"

    def test_copies_defRPr(self):
        """source has defRPr with font size → copy → target has identical defRPr."""
        source_p = _make_paragraph()
        source_pPr = source_p.find(qn("a:pPr"))
        defRPr = etree.SubElement(source_pPr, qn("a:defRPr"))
        defRPr.set("sz", "1400")

        target_p = _make_paragraph(with_pPr=False)

        copy_paragraph_properties(source_pPr, target_p)

        target_pPr = target_p.find(qn("a:pPr"))
        defRPr_copy = target_pPr.find(qn("a:defRPr"))
        assert defRPr_copy is not None
        assert defRPr_copy.get("sz") == "1400"

    def test_handles_empty_source_pPr(self):
        """source pPr is empty element → copy → target gets empty pPr (no crash)."""
        source_p = _make_paragraph()
        source_pPr = source_p.find(qn("a:pPr"))
        # Remove all children from source pPr
        for child in list(source_pPr):
            source_pPr.remove(child)

        target_p = _make_paragraph(with_pPr=False)

        copy_paragraph_properties(source_pPr, target_p)

        target_pPr = target_p.find(qn("a:pPr"))
        assert target_pPr is not None
        assert len(target_pPr) == 0  # Empty pPr

    def test_handles_none_source(self):
        """source_pPr is None → function returns without error (defensive)."""
        target_p = _make_paragraph(with_pPr=False)

        # Should not raise
        copy_paragraph_properties(None, target_p)


class TestCopyRunProperties:
    """Tests for copy_run_properties."""

    def test_copies_rPr(self):
        """source run has rPr with bold, size → copy → target run has matching rPr."""
        source_r = etree.Element(qn("a:r"))
        source_rPr = etree.SubElement(source_r, qn("a:rPr"))
        source_rPr.set("b", "1")
        source_rPr.set("sz", "1400")

        target_r = etree.Element(qn("a:r"))

        copy_run_properties(source_r, target_r)

        target_rPr = target_r.find(qn("a:rPr"))
        assert target_rPr is not None
        assert target_rPr.get("b") == "1"
        assert target_rPr.get("sz") == "1400"

    def test_handles_missing_source_rPr(self):
        """source run has no rPr → function does nothing."""
        source_r = etree.Element(qn("a:r"))
        target_r = etree.Element(qn("a:r"))

        # Should not raise
        copy_run_properties(source_r, target_r)

        target_rPr = target_r.find(qn("a:rPr"))
        assert target_rPr is None

    def test_replaces_existing_target_rPr(self):
        """target has rPr → copy replaces it with source's rPr."""
        source_r = etree.Element(qn("a:r"))
        source_rPr = etree.SubElement(source_r, qn("a:rPr"))
        source_rPr.set("sz", "1400")

        target_r = etree.Element(qn("a:r"))
        target_rPr = etree.SubElement(target_r, qn("a:rPr"))
        target_rPr.set("sz", "1200")

        copy_run_properties(source_r, target_r)

        target_rPr = target_r.find(qn("a:rPr"))
        assert target_rPr.get("sz") == "1400"

    def test_deep_copy_isolation(self):
        """modify target rPr → source unchanged."""
        source_r = etree.Element(qn("a:r"))
        source_rPr = etree.SubElement(source_r, qn("a:rPr"))
        source_rPr.set("sz", "1400")

        target_r = etree.Element(qn("a:r"))

        copy_run_properties(source_r, target_r)

        # Modify target
        target_rPr = target_r.find(qn("a:rPr"))
        target_rPr.set("sz", "9999")

        # Source should be unchanged
        assert source_rPr.get("sz") == "1400"


class TestApplyBulletType:
    """Tests for apply_bullet_type."""

    def test_bullet_char_explicit(self):
        """BulletSpec(type='bullet', char='–') → p gets buChar with char='–'."""
        p = _make_paragraph(with_pPr=False)
        spec = BulletSpec(type="bullet", char="–")

        apply_bullet_type(p, spec)

        pPr = p.find(qn("a:pPr"))
        buChar = pPr.find(qn("a:buChar"))
        assert buChar is not None
        assert buChar.get("char") == "–"

    def test_bullet_char_from_template(self):
        """BulletSpec(type='bullet') + template_pPr has buChar='▸' → p gets buChar='▸'."""
        p = _make_paragraph(with_pPr=False)
        spec = BulletSpec(type="bullet")

        template_p = _make_paragraph(bullet_char="▸")
        template_pPr = template_p.find(qn("a:pPr"))

        apply_bullet_type(p, spec, template_pPr)

        pPr = p.find(qn("a:pPr"))
        buChar = pPr.find(qn("a:buChar"))
        assert buChar is not None
        assert buChar.get("char") == "▸"

    def test_bullet_char_default_fallback(self):
        """BulletSpec(type='bullet') + no template → p gets buChar='•'."""
        p = _make_paragraph(with_pPr=False)
        spec = BulletSpec(type="bullet")

        apply_bullet_type(p, spec)

        pPr = p.find(qn("a:pPr"))
        buChar = pPr.find(qn("a:buChar"))
        assert buChar is not None
        assert buChar.get("char") == "•"

    def test_number_default_scheme(self):
        """BulletSpec(type='number') → p gets buAutoNum type='arabicPeriod'."""
        p = _make_paragraph(with_pPr=False)
        spec = BulletSpec(type="number")

        apply_bullet_type(p, spec)

        pPr = p.find(qn("a:pPr"))
        buAutoNum = pPr.find(qn("a:buAutoNum"))
        assert buAutoNum is not None
        assert buAutoNum.get("type") == "arabicPeriod"

    def test_number_custom_scheme(self):
        """BulletSpec(type='number', scheme='alphaLcPeriod') → p gets buAutoNum type='alphaLcPeriod'."""
        p = _make_paragraph(with_pPr=False)
        spec = BulletSpec(type="number", scheme="alphaLcPeriod")

        apply_bullet_type(p, spec)

        pPr = p.find(qn("a:pPr"))
        buAutoNum = pPr.find(qn("a:buAutoNum"))
        assert buAutoNum is not None
        assert buAutoNum.get("type") == "alphaLcPeriod"

    def test_number_start_at(self):
        """BulletSpec(type='number', start_at=3) → p gets buAutoNum startAt='3'."""
        p = _make_paragraph(with_pPr=False)
        spec = BulletSpec(type="number", start_at=3)

        apply_bullet_type(p, spec)

        pPr = p.find(qn("a:pPr"))
        buAutoNum = pPr.find(qn("a:buAutoNum"))
        assert buAutoNum is not None
        assert buAutoNum.get("startAt") == "3"

    def test_none_removes_bullets(self):
        """BulletSpec(type='none') → p gets buNone, existing buChar removed."""
        p = _make_paragraph(bullet_char="•")
        spec = BulletSpec(type="none")

        apply_bullet_type(p, spec)

        pPr = p.find(qn("a:pPr"))
        buNone = pPr.find(qn("a:buNone"))
        assert buNone is not None
        # buChar should be gone
        buChar = pPr.find(qn("a:buChar"))
        assert buChar is None

    def test_auto_copies_from_template(self):
        """BulletSpec(type='auto') + template has buChar → p gets buChar."""
        p = _make_paragraph(with_pPr=False)
        spec = BulletSpec(type="auto")

        template_p = _make_paragraph(bullet_char="▸")
        template_pPr = template_p.find(qn("a:pPr"))

        apply_bullet_type(p, spec, template_pPr)

        pPr = p.find(qn("a:pPr"))
        buChar = pPr.find(qn("a:buChar"))
        assert buChar is not None
        assert buChar.get("char") == "▸"

    def test_auto_no_template_noop(self):
        """BulletSpec(type='auto') + no template → p unchanged."""
        p = _make_paragraph(with_pPr=False)
        spec = BulletSpec(type="auto")

        apply_bullet_type(p, spec)

        pPr = p.find(qn("a:pPr"))
        # Should just have empty pPr
        assert len(pPr) == 0

    def test_removes_existing_before_applying(self):
        """p has buChar → apply number → buChar gone, buAutoNum present."""
        p = _make_paragraph(bullet_char="•")
        spec = BulletSpec(type="number")

        apply_bullet_type(p, spec)

        pPr = p.find(qn("a:pPr"))
        buChar = pPr.find(qn("a:buChar"))
        assert buChar is None
        buAutoNum = pPr.find(qn("a:buAutoNum"))
        assert buAutoNum is not None

    def test_mutually_exclusive_removal(self):
        """p has buAutoNum → apply buNone → buAutoNum gone, buNone present."""
        p = _make_paragraph(autonum="arabicPeriod")
        spec = BulletSpec(type="none")

        apply_bullet_type(p, spec)

        pPr = p.find(qn("a:pPr"))
        buAutoNum = pPr.find(qn("a:buAutoNum"))
        assert buAutoNum is None
        buNone = pPr.find(qn("a:buNone"))
        assert buNone is not None

    def test_element_ordering_after_apply(self):
        """p has defRPr → apply buChar → buChar is before defRPr in XML."""
        p = _make_paragraph()
        pPr = p.find(qn("a:pPr"))
        etree.SubElement(pPr, qn("a:defRPr"))

        spec = BulletSpec(type="bullet", char="•")

        apply_bullet_type(p, spec)

        tags = _get_child_tags(pPr)
        buChar_idx = tags.index("buChar")
        defRPr_idx = tags.index("defRPr")
        assert buChar_idx < defRPr_idx

    def test_ensures_pPr_if_missing(self):
        """p has no pPr → apply bullet → pPr created, bullet added correctly."""
        p = _make_paragraph(with_pPr=False)
        spec = BulletSpec(type="bullet", char="•")

        apply_bullet_type(p, spec)

        pPr = p.find(qn("a:pPr"))
        assert pPr is not None
        buChar = pPr.find(qn("a:buChar"))
        assert buChar is not None


class TestResolveLevelStyle:
    """Tests for resolve_level_style."""

    def test_resolves_level_0_from_lstStyle(self):
        """placeholder with lstStyle containing lvl1pPr → resolves level 0 style."""
        # Create mock placeholder element
        ph = etree.Element(qn("p:sp"))
        txBody = etree.SubElement(ph, qn("p:txBody"))
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        lvl1pPr = etree.SubElement(lstStyle, qn("a:lvl1pPr"))
        lvl1pPr.set("marL", "457200")

        # Create mock placeholder object
        class MockPlaceholder:
            def __init__(self, elem):
                self._element = elem

        placeholder = MockPlaceholder(ph)

        styles = resolve_level_style(placeholder, 0)

        assert "indent" in styles
        assert styles["indent"] == 457200

    def test_resolves_level_1_from_lstStyle(self):
        """placeholder with lvl2pPr → resolves level 1 (lstStyle uses 1-indexed names)."""
        ph = etree.Element(qn("p:sp"))
        txBody = etree.SubElement(ph, qn("p:txBody"))
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        lvl2pPr = etree.SubElement(lstStyle, qn("a:lvl2pPr"))
        lvl2pPr.set("marL", "685800")

        class MockPlaceholder:
            def __init__(self, elem):
                self._element = elem

        placeholder = MockPlaceholder(ph)

        styles = resolve_level_style(placeholder, 1)

        assert "indent" in styles
        assert styles["indent"] == 685800

    def test_returns_empty_for_undefined_level(self):
        """placeholder with only lvl1pPr → resolve level 5 → empty dict."""
        ph = etree.Element(qn("p:sp"))
        txBody = etree.SubElement(ph, qn("p:txBody"))
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        etree.SubElement(lstStyle, qn("a:lvl1pPr"))

        class MockPlaceholder:
            def __init__(self, elem):
                self._element = elem

        placeholder = MockPlaceholder(ph)

        styles = resolve_level_style(placeholder, 5)

        assert styles == {}

    def test_returns_empty_for_no_lstStyle(self):
        """placeholder without lstStyle → empty dict."""
        ph = etree.Element(qn("p:sp"))
        etree.SubElement(ph, qn("p:txBody"))

        class MockPlaceholder:
            def __init__(self, elem):
                self._element = elem

        placeholder = MockPlaceholder(ph)

        styles = resolve_level_style(placeholder, 0)

        assert styles == {}

    def test_extracts_spacing_as_child_elements(self):
        """lstStyle lvl1pPr has spcBef>spcPts val='600' → space_before = 6.0 (points)."""
        ph = etree.Element(qn("p:sp"))
        txBody = etree.SubElement(ph, qn("p:txBody"))
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        lvl1pPr = etree.SubElement(lstStyle, qn("a:lvl1pPr"))

        spcBef = etree.SubElement(lvl1pPr, qn("a:spcBef"))
        spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
        spcPts.set("val", "600")

        class MockPlaceholder:
            def __init__(self, elem):
                self._element = elem

        placeholder = MockPlaceholder(ph)

        styles = resolve_level_style(placeholder, 0)

        assert "space_before" in styles
        assert styles["space_before"] == 6.0

    def test_extracts_line_spacing_pct(self):
        """lstStyle lvl1pPr has lnSpc>spcPct val='100000' → line_spacing = 100000."""
        ph = etree.Element(qn("p:sp"))
        txBody = etree.SubElement(ph, qn("p:txBody"))
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        lvl1pPr = etree.SubElement(lstStyle, qn("a:lvl1pPr"))

        lnSpc = etree.SubElement(lvl1pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", "100000")

        class MockPlaceholder:
            def __init__(self, elem):
                self._element = elem

        placeholder = MockPlaceholder(ph)

        styles = resolve_level_style(placeholder, 0)

        assert "line_spacing" in styles
        assert styles["line_spacing"] == 100000

    def test_extracts_bullet_type_char(self):
        """lstStyle lvl1pPr has buChar char='•' → bullet_type = 'char', bullet_char = '•'."""
        ph = etree.Element(qn("p:sp"))
        txBody = etree.SubElement(ph, qn("p:txBody"))
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        lvl1pPr = etree.SubElement(lstStyle, qn("a:lvl1pPr"))

        buChar = etree.SubElement(lvl1pPr, qn("a:buChar"))
        buChar.set("char", "•")

        class MockPlaceholder:
            def __init__(self, elem):
                self._element = elem

        placeholder = MockPlaceholder(ph)

        styles = resolve_level_style(placeholder, 0)

        assert styles["bullet_type"] == "char"
        assert styles["bullet_char"] == "•"

    def test_extracts_bullet_type_autonum(self):
        """lstStyle lvl1pPr has buAutoNum → bullet_type = 'autonum'."""
        ph = etree.Element(qn("p:sp"))
        txBody = etree.SubElement(ph, qn("p:txBody"))
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        lvl1pPr = etree.SubElement(lstStyle, qn("a:lvl1pPr"))

        buAutoNum = etree.SubElement(lvl1pPr, qn("a:buAutoNum"))
        buAutoNum.set("type", "arabicPeriod")

        class MockPlaceholder:
            def __init__(self, elem):
                self._element = elem

        placeholder = MockPlaceholder(ph)

        styles = resolve_level_style(placeholder, 0)

        assert styles["bullet_type"] == "autonum"


# ============================================================================
# Integration Tests with Real Template
# ============================================================================


@pytest.fixture
def real_template():
    """Get the Uptimize Master template."""
    from pptx_mcp.template_engine import TemplateEngine

    template_path = Path(__file__).parent.parent / "templates"
    templates = list(template_path.glob("*.potx")) + list(template_path.glob("*.pptx"))
    if not templates:
        pytest.skip("No template found for integration test")

    # Use template_engine to handle .potx conversion
    engine = TemplateEngine(str(template_path))
    template_file = templates[0]
    try:
        info = engine.register_template(template_file)
        prs = engine.open_presentation(info.name)
        # Add a slide so we have something to work with
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])
        return prs
    finally:
        engine.cleanup()


@pytest.fixture
def real_placeholder(real_template):
    """Get a content placeholder from the real template."""
    # Find a layout with a content placeholder
    for layout in real_template.slide_layouts:
        for ph in layout.placeholders:
            try:
                ph_type = str(ph.placeholder_format.type)
                if "OBJECT" in ph_type or "BODY" in ph_type:
                    return ph
            except Exception:
                continue
    pytest.skip("No content placeholder found in template")


class TestRealTemplate:
    """Integration tests with real PowerPoint template."""

    def test_copy_real_template_pPr(self, real_placeholder):
        """Get first paragraph from real placeholder, copy pPr to new paragraph."""
        if not hasattr(real_placeholder, "text_frame"):
            pytest.skip("Placeholder has no text frame")

        tf = real_placeholder.text_frame
        if not tf.paragraphs:
            pytest.skip("Text frame has no paragraphs")

        first_p = tf.paragraphs[0]
        source_pPr = first_p._element.find(qn("a:pPr"))
        if source_pPr is None:
            pytest.skip("First paragraph has no pPr")

        # Create a new paragraph element
        target_p = _make_paragraph(with_pPr=False)

        copy_paragraph_properties(source_pPr, target_p)

        target_pPr = target_p.find(qn("a:pPr"))
        assert target_pPr is not None
        # Verify copy has similar structure (at least some children)
        # Note: exact structure depends on template
        assert len(target_pPr) > 0 or len(source_pPr.attrib) > 0

    def test_resolve_real_template_level_0(self, real_placeholder):
        """Call resolve_level_style on real placeholder at level 0."""
        styles = resolve_level_style(real_placeholder, 0)

        # Should return non-empty dict if template has lstStyle
        # (Uptimize Master is known to have lstStyle)
        # Even if empty, should not crash
        assert isinstance(styles, dict)
